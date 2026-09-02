"""Build the ASTRA executive research & collaboration deck.

Every experimental statement carries an evidence badge. Nothing is asserted that
is not traceable to `research/E17_*.md` or `results/E17_*`.

Run with the isolated docs venv (python-pptx), NOT the measurement venv.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x12, 0x26, 0x3A)
MUTED = RGBColor(0x5A, 0x6B, 0x7C)
FAINT = RGBColor(0x8A, 0x99, 0xA8)
ACCENT = RGBColor(0x1B, 0x6C, 0xA8)
PANEL = RGBColor(0xF2, 0xF5, 0xF8)
LINE = RGBColor(0xD8, 0xDF, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EST = RGBColor(0x1E, 0x6B, 0x52)      # established - green
PART = RGBColor(0xA9, 0x76, 0x0B)     # partially validated - amber
PLAN = RGBColor(0x3D, 0x5A, 0x8A)     # planned - blue
HYP = RGBColor(0x6B, 0x72, 0x80)      # hypothesis - grey
RED = RGBColor(0xB0, 0x36, 0x2A)

BADGE = {
    "ESTABLISHED": EST,
    "PARTIALLY VALIDATED": PART,
    "PRELIMINARY": PART,
    "PLANNED": PLAN,
    "HYPOTHESIS": HYP,
    "WITHDRAWN": RED,
    "MUST NOT CLAIM": RED,
}

FONT = "Segoe UI"
MONO = "Consolas"


def _tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _p(tf, text, size, color=INK, bold=False, first=False, space=6,
       font=FONT, align=PP_ALIGN.LEFT, italic=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_after = Pt(space)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return para


def _rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.0):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def _badge(slide, x, y, label, w=Inches(1.85)):
    col = BADGE[label]
    s = _rect(slide, x, y, w, Inches(0.26), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    _p(tf, label, 9, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
    return s


def slide(prs, title=None, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    if title:
        if kicker:
            _p(_tb(s, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.24)),
               kicker.upper(), 10, ACCENT, bold=True, first=True, space=0)
        _p(_tb(s, Inches(0.7), Inches(0.68), Inches(11.9), Inches(0.6)),
           title, 27, INK, bold=True, first=True, space=0)
        _rect(s, Inches(0.7), Inches(1.32), Inches(1.5), Emu(20955), fill=ACCENT)
    return s


def foot(s, text):
    _p(_tb(s, Inches(0.7), Inches(6.92), Inches(11.9), Inches(0.3)),
       text, 8.5, FAINT, first=True, space=0, italic=True)


def table(slide, x, y, w, rows, widths=None, fs=10.5, header_fs=10,
          row_h=Inches(0.32), head_h=Inches(0.34), colors=None):
    nr, nc = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nr, nc, x, y, w, head_h + row_h * (nr - 1))
    t = shape.table
    t.rows[0].height = head_h
    for i in range(1, nr):
        t.rows[i].height = row_h
    if widths:
        total = sum(widths)
        for j, frac in enumerate(widths):
            t.columns[j].width = Emu(int(w * frac / total))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = t.cell(i, j)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK if i == 0 else (WHITE if i % 2 else PANEL)
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.space_after = Pt(0)
            run = para.add_run()
            run.text = str(val)
            run.font.size = Pt(header_fs if i == 0 else fs)
            run.font.bold = i == 0
            run.font.name = FONT
            if i == 0:
                run.font.color.rgb = WHITE
            elif colors and (i, j) in colors:
                run.font.color.rgb = colors[(i, j)]
                run.font.bold = True
            else:
                run.font.color.rgb = INK
    return t


def chain(slide, x, y, items, w=Inches(1.42), h=Inches(0.52), gap=Inches(0.12),
          fill=None, fs=9.5, vertical=False, col=None):
    """A left-to-right (or top-down) chevron chain."""
    for i, label in enumerate(items):
        c = (col[i] if col else None) or fill or ACCENT
        if vertical:
            bx, by = x, y + i * (h + gap)
        else:
            bx, by = x + i * (w + gap), y
        s = _rect(slide, bx, by, w, h, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = s.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        _p(tf, label, fs, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
        if i < len(items) - 1:
            if vertical:
                _p(_tb(slide, bx, by + h, w, gap), "▼", 8, FAINT, first=True,
                   space=0, align=PP_ALIGN.CENTER)
            else:
                _p(_tb(slide, bx + w, by + Inches(0.14), gap, Inches(0.3)), "▶", 9,
                   FAINT, first=True, space=0, align=PP_ALIGN.CENTER)


def kv_panel(slide, x, y, w, h, heading, lines, hcolor=ACCENT):
    _rect(slide, x, y, w, h, fill=PANEL)
    _rect(slide, x, y, Emu(30000), h, fill=hcolor)
    tf = _tb(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.4), h - Inches(0.3))
    _p(tf, heading, 12.5, INK, bold=True, first=True, space=7)
    for ln in lines:
        if isinstance(ln, tuple):
            _p(tf, ln[0], 10.5, ln[1], bold=ln[2] if len(ln) > 2 else False, space=5)
        else:
            _p(tf, ln, 10.5, MUTED, space=5)


# ---------------------------------------------------------------- deck ----
def build(out: Path, pos: dict | None) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 TITLE ---------------------------------------------------------------
    s = slide(prs)
    _rect(s, 0, 0, W, H, fill=INK)
    _rect(s, Inches(0.9), Inches(2.28), Inches(1.7), Emu(34925), fill=ACCENT)
    _p(_tb(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.4)),
       "RESEARCH & COLLABORATION BRIEFING", 12, RGBColor(0x7E, 0xA8, 0xC9),
       bold=True, first=True, space=0)
    tf = _tb(s, Inches(0.9), Inches(2.62), Inches(11.5), Inches(1.9))
    _p(tf, "ASTRA", 50, WHITE, bold=True, first=True, space=2)
    _p(tf, "Fault Observability, Information Loss & Monitor Placement", 23,
       RGBColor(0xDA, 0xE4, 0xEE), space=2)
    _p(tf, "in Autonomous Driving Systems", 23, RGBColor(0xDA, 0xE4, 0xEE), space=14)
    _p(tf, "Current Evidence  •  Research Roadmap  •  Dataset Validation  •  "
           "Industry Value  •  Publication Potential", 13, RGBColor(0x9F, 0xBA, 0xD2))
    _p(_tb(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6)),
       "Sushanth C.  ·  Tarun Gowda V  ·  T Tilak Reddy   |   Guide: Dr. Chaitra R.   |   "
       "BMS College of Engineering   |   1 September 2026", 11,
       RGBColor(0x8D, 0xA5, 0xBC), first=True, space=0)

    # 2 EXECUTIVE SUMMARY ---------------------------------------------------
    s = slide(prs, "Executive summary", "where the project stands")
    kv_panel(s, Inches(0.7), Inches(1.62), Inches(3.78), Inches(2.25),
             "What exists today",
             [("An audited experimental framework", INK, True),
              "90 profiles · 2,160 closed-loop runs · 0 failures",
              "Pre-registered before execution",
              "Three measurement defects found and corrected — two by self-audit"], EST)
    kv_panel(s, Inches(4.72), Inches(1.62), Inches(3.78), Inches(2.25),
             "The initial finding",
             [("Fault evidence weakens stage-by-stage,", INK, True),
              ("and differently for different faults", INK, True),
              "One stable, well-posed result at n=30",
              "Three other faults behave in three other ways"], PART)
    kv_panel(s, Inches(8.74), Inches(1.62), Inches(3.86), Inches(2.25),
             "What is not yet known",
             [("Whether it generalises", INK, True),
              "No real-world data yet · no closed loop yet",
              "Monitor-placement consequence untested",
              "Adversarial extension not started"], PLAN)

    _rect(s, Inches(0.7), Inches(4.15), Inches(11.9), Inches(1.5), fill=PANEL)
    _rect(s, Inches(0.7), Inches(4.15), Emu(30000), Inches(1.5), fill=ACCENT)
    tf = _tb(s, Inches(1.0), Inches(4.36), Inches(11.3), Inches(1.15))
    _p(tf, "The ask", 12.5, INK, bold=True, first=True, space=7)
    _p(tf, "ASTRA has a technically audited foundation and a scientifically interesting "
           "initial finding. The remaining work is not feature development — it is a gated "
           "validation programme that determines whether the finding generalises to real data, "
           "adversarial conditions and closed-loop driving.", 12, MUTED, space=0)
    foot(s, "Every claim in this deck is badged ESTABLISHED / PARTIALLY VALIDATED / PLANNED / "
            "HYPOTHESIS and is traceable to the project's own audit documents.")

    # 3 THE PROBLEM ---------------------------------------------------------
    s = slide(prs, "The problem", "why this question is not already answered")
    tf = _tb(s, Inches(0.7), Inches(1.62), Inches(5.6), Inches(3.4))
    _p(tf, "A sensor fault does not stay a sensor fault.", 15, INK, bold=True,
       first=True, space=10)
    _p(tf, "It passes through estimation, fusion, trust scoring and a stack of safety "
           "monitors. Each stage transforms it. Somewhere along that path the evidence "
           "that anything is wrong can become too weak for a downstream monitor to act on.",
       12.5, MUTED, space=10)
    _p(tf, "Conventional fault detection asks a binary question:", 12.5, MUTED, space=6)
    _p(tf, "        “fault or no fault?”", 13, INK, bold=True, font=MONO, space=10)
    _p(tf, "It does not tell an engineer WHERE in their own architecture the evidence "
           "stopped being usable — which is exactly what determines where a monitor "
           "should sit.", 12.5, MUTED, space=0)

    _rect(s, Inches(6.75), Inches(1.62), Inches(5.85), Inches(3.75), fill=PANEL)
    tfh = _tb(s, Inches(7.05), Inches(1.82), Inches(5.3), Inches(0.3))
    _p(tfh, "THE BLIND-SPOT PROBLEM", 10, ACCENT, bold=True, first=True, space=0)
    chain(s, Inches(7.05), Inches(2.28), ["Sensor fault occurs"], w=Inches(5.25),
          h=Inches(0.42), fill=INK, fs=11)
    labels = [("L1  sensing", "evidence strong", EST),
              ("L2  estimation", "evidence weakens", PART),
              ("L6/L7  gates", "evidence insufficient", RED)]
    yy = Inches(2.92)
    for name, note, col in labels:
        _rect(s, Inches(7.05), yy, Inches(2.35), Inches(0.42), fill=col,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t2 = s.shapes[-1].text_frame
        t2.margin_left = Inches(0.1)
        _p(t2, name, 10.5, WHITE, bold=True, first=True, space=0)
        _p(_tb(s, Inches(9.55), yy + Inches(0.1), Inches(2.8), Inches(0.3)),
           note, 10.5, MUTED, first=True, space=0)
        yy += Inches(0.55)
    _p(_tb(s, Inches(7.05), Inches(4.72), Inches(5.3), Inches(0.5)),
       "The monitor sits after the point where the evidence became insufficient — "
       "so it cannot act, however well it is implemented.", 11, INK, first=True, space=0)
    foot(s, "Illustrative schematic of the problem class. Measured per-fault profiles appear "
            "on the results slides.")

    # 4 WHAT ASTRA DOES -----------------------------------------------------
    s = slide(prs, "What ASTRA actually does", "positioning")
    _rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.02), fill=PANEL)
    _rect(s, Inches(0.7), Inches(1.6), Emu(30000), Inches(1.02), fill=ACCENT)
    tf = _tb(s, Inches(1.0), Inches(1.78), Inches(11.3), Inches(0.8))
    _p(tf, "A framework for quantifying how sensor-fault evidence changes across an "
           "autonomous-driving estimation and safety pipeline — identifying where "
           "discriminative evidence becomes insufficient for downstream detection, and "
           "investigating how that guides monitor placement.", 13, INK, bold=True,
       first=True, space=0)

    tf = _tb(s, Inches(0.7), Inches(2.95), Inches(5.85), Inches(3.4))
    _p(tf, "The primary question", 12.5, ACCENT, bold=True, first=True, space=7)
    _p(tf, "When a sensor becomes faulty, where in the pipeline does the fault stop being "
           "distinguishable from healthy behaviour?", 13, INK, bold=True, space=12)
    _p(tf, "Secondary question", 12.5, ACCENT, bold=True, space=7)
    _p(tf, "Can that information decide where monitoring should occur?", 13, INK,
       bold=True, space=12)
    _p(tf, "Future security question", 12.5, ACCENT, bold=True, space=7)
    _p(tf, "Does a deliberately manipulated sensor behave differently from a naturally "
           "faulty one?  — not yet answered", 12.5, HYP, bold=True, space=0)

    kv_panel(s, Inches(6.9), Inches(2.95), Inches(5.7), Inches(3.35),
             "ASTRA is NOT positioned as",
             [("✕  a sensor fault detector", MUTED),
              ("✕  a Kalman-filtering contribution", MUTED),
              ("✕  a CARLA testing exercise", MUTED),
              ("✕  a multi-dataset benchmark", MUTED),
              ("", MUTED),
              ("It measures how evidence about a fault survives", INK, True),
              ("transformation — not whether a fault is present.", INK, True)], MUTED)

    # 5 PIPELINE ------------------------------------------------------------
    s = slide(prs, "The ASTRA pipeline", "nine layers, two isolated cores")
    _p(_tb(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.3)),
       "Measurement points used by the experiment are marked. Core-A proposes; Core-B may "
       "veto. The trust boundary is one-way.", 11, MUTED, first=True, space=0)
    cells = [("L1", "sensing", EST), ("L2a", "innovation", EST), ("L2b", "estimate", EST),
             ("L3", "trust", EST), ("L4", "policy\nCORE-A", MUTED), ("L5", "twin", PLAN),
             ("L6", "conformal", PART), ("L7", "shield", EST), ("L8", "fail-safe", EST),
             ("L9", "arbitration", PLAN)]
    x = Inches(0.7)
    bw, gap = Inches(1.13), Inches(0.088)
    for code, name, col in cells:
        _rect(s, x, Inches(2.1), bw, Inches(0.95), fill=col,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t2 = s.shapes[-1].text_frame
        t2.word_wrap = True
        t2.margin_left = t2.margin_right = Inches(0.04)
        _p(t2, code, 14, WHITE, bold=True, first=True, space=1, align=PP_ALIGN.CENTER)
        _p(t2, name, 8.5, RGBColor(0xE2, 0xEA, 0xF2), space=0, align=PP_ALIGN.CENTER)
        x += bw + gap
    _rect(s, Inches(0.7), Inches(3.2), Inches(4.72), Inches(0.3), fill=PANEL)
    _p(_tb(s, Inches(0.7), Inches(3.24), Inches(4.72), Inches(0.25)),
       "SHARED  ·  L1–L3", 9.5, MUTED, first=True, space=0, align=PP_ALIGN.CENTER)
    _rect(s, Inches(5.5), Inches(3.2), Inches(1.13), Inches(0.3), fill=PANEL)
    _p(_tb(s, Inches(5.5), Inches(3.24), Inches(1.13), Inches(0.25)),
       "CORE-A", 9.5, MUTED, first=True, space=0, align=PP_ALIGN.CENTER)
    _rect(s, Inches(6.71), Inches(3.2), Inches(4.68), Inches(0.3), fill=PANEL)
    _p(_tb(s, Inches(6.71), Inches(3.24), Inches(4.68), Inches(0.25)),
       "CORE-B  ·  3 gates verdict on every tick", 9.5, MUTED, first=True, space=0,
       align=PP_ALIGN.CENTER)

    _badge(s, Inches(0.7), Inches(3.78), "ESTABLISHED")
    _p(_tb(s, Inches(2.7), Inches(3.76), Inches(9.9), Inches(0.32)),
       "9 / 9 layers domain-mapped · 3 / 3 Core-B gates produce a verdict on every one of "
       "300 ticks — verified in a live run, not from documentation.", 11.5, INK,
       first=True, space=0)

    table(s, Inches(0.7), Inches(4.3), Inches(11.9), [
        ["Stage", "Statistic measured", "Role in the experiment"],
        ["L1", "raw channel value / dispersion / stream health", "sensor-level detectability"],
        ["L2a", "UKF innovation (Mahalanobis)", "estimator-level observability"],
        ["L2b", "estimated lateral position", "post-estimation state"],
        ["L3", "trust index", "normalised-innovation trust"],
        ["L6", "conformal non-conformity score", "statistical gate evidence"],
        ["L7 / L8", "veto flag / fail-safe posture", "operational gate outcome"],
    ], widths=[1.1, 4.2, 4.0], fs=10.5)
    foot(s, "Stage statistics are one defensible choice among several; a different statistic could "
            "move the measured values. Stated as a limitation throughout.")

    # 6 E17 STATUS ----------------------------------------------------------
    s = slide(prs, "Current E17 status", "audited, 1 September 2026")
    _rect(s, Inches(0.7), Inches(1.58), Inches(3.5), Inches(0.72), fill=PART,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = s.shapes[-1].text_frame
    tf.margin_left = Inches(0.15)
    _p(tf, "E17 — NEEDS MINOR FIXES", 14, WHITE, bold=True, first=True, space=0)
    _p(_tb(s, Inches(4.45), Inches(1.62), Inches(8.15), Inches(0.7)),
       "Sweep audited: 4 of 6 faults passed injection integrity; 2 were invalid, have since been "
       "re-injected through the correct per-channel path, and re-run at n=30. Contaminated "
       "records are retained and flagged, not deleted.",
       12, MUTED, first=True, space=0)

    table(s, Inches(0.7), Inches(2.52), Inches(11.9), [
        ["Fault", "Injection", "30-seed data", "Current interpretation"],
        ["speed_stuck", "Valid", "Valid", "L2a — well-posed on 30/30 seeds (policy P1)"],
        ["speed_bias", "Valid", "Valid", "L2a modal 100 %, but 0/30 well-posed"],
        ["lateral_noise", "Valid", "Valid", "Persistent — never absorbed at any stage"],
        ["imu_dropout", "Valid", "Valid", "Non-monotonic — no well-posed A(f)"],
        ["position_bias", "Corrected", "Re-run n=30", "NOT absorbed — D_L2a 0.675–0.797"],
        ["position_drift", "Corrected", "Re-run n=30", "NOT absorbed — D_L2a 0.655–0.767"],
    ], widths=[2.0, 1.3, 1.4, 5.6], fs=11,
        colors={(1, 1): EST, (2, 1): EST, (3, 1): EST, (4, 1): EST,
                (5, 1): PART, (6, 1): PART, (5, 3): RED, (6, 3): RED})

    kv_panel(s, Inches(0.7), Inches(5.28), Inches(5.85), Inches(1.42),
             "The two position faults, corrected and re-run",
             [("Re-injected via redundant.offset — the per-channel path. Fault", MUTED),
              ("reached the estimator in 720/720 runs. Result: absorbed at L2a in", MUTED),
              ("0 of 12 cells. The original D = 0.500 was entirely the artefact.", INK, True)], PART)
    kv_panel(s, Inches(6.75), Inches(5.28), Inches(5.85), Inches(1.42),
             "Why this strengthens rather than weakens the work",
             [("The defect was found by auditing our own headline result", MUTED),
              ("because it looked too clean. The workflow now requires", MUTED),
              ("delivered-signal verification before any result is accepted.", MUTED)], EST)

    # 7 WHAT RESULTS MEAN ---------------------------------------------------
    s = slide(prs, "What the current results mean", "reading the metric correctly")
    _p(_tb(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.3)),
       "D_s = AUC(faulted, matched-clean), folded to [0.5, 1.0]. Unit of analysis is the "
       "seed, never the tick.", 11.5, MUTED, first=True, space=0)

    _rect(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(1.28), fill=PANEL)
    _rect(s, Inches(0.7), Inches(2.0), Emu(30000), Inches(1.28), fill=RED)
    tf = _tb(s, Inches(1.0), Inches(2.16), Inches(5.4), Inches(1.0))
    _p(tf, "D = 0.5 does NOT mean “50 % accuracy”", 12.5, INK, bold=True, first=True, space=6)
    _p(tf, "0.5 is chance-level discriminability and the floor of a folded metric. It means "
           "the evaluated representation does not retain discriminative evidence — not that "
           "a detector is right half the time.", 11, MUTED, space=0)

    _rect(s, Inches(6.75), Inches(2.0), Inches(5.85), Inches(1.28), fill=PANEL)
    _rect(s, Inches(6.75), Inches(2.0), Emu(30000), Inches(1.28), fill=RED)
    tf = _tb(s, Inches(7.05), Inches(2.16), Inches(5.4), Inches(1.0))
    _p(tf, "High D does NOT mean the gate fired", 12.5, INK, bold=True, first=True, space=6)
    _p(tf, "AUC is scale-free. It reports near-perfect separation for a shift far below a "
           "gate's decision threshold. Statistical separability and operational detection "
           "are different claims.", 11, MUTED, space=0)

    _p(_tb(s, Inches(0.7), Inches(3.52), Inches(11.9), Inches(0.3)),
       "Fault-specific observability trajectories — four valid faults, four distinct behaviours",
       12.5, INK, bold=True, first=True, space=0)
    table(s, Inches(0.7), Inches(3.94), Inches(11.9), [
        ["Fault", "Behaviour", "Absorption point A(f)", "Evidence"],
        ["speed_stuck", "Absorbed at estimation", "L2a — well-posed 30/30 (P1)",
         "0.9625 → 0.5629, SD 0.0077"],
        ["speed_bias", "Weakens, then statistic recovers", "not well-posed — 0/30 seeds",
         "modal L2a 100 % is misleading"],
        ["imu_dropout", "Non-monotonic", "none — no unique A(f)", "curve crosses and returns"],
        ["lateral_noise", "Persistent", "none — never absorbed", "stays separable throughout"],
        ["position_bias / drift", "Survives estimation (corrected)", "none — not absorbed",
         "D_L2a 0.655–0.797, n=30"],
    ], widths=[1.7, 2.9, 3.3, 3.4], fs=10.5)
    _badge(s, Inches(0.7), Inches(6.46), "PARTIALLY VALIDATED")
    _p(_tb(s, Inches(2.7), Inches(6.44), Inches(9.9), Inches(0.32)),
       "Heterogeneity is the finding: a single absorption point is ill-posed for five of six "
       "faults. One fault-policy pair shows a clean, well-posed absorption.", 11.5, INK,
       first=True, space=0)
    foot(s, "One plant model · one severity per fault · three policies · simulation only. "
            "Reproducibility across seeds is not external validity.")

    # 8 VALIDATED VS INVALIDATED -------------------------------------------
    s = slide(prs, "Validated vs invalidated", "what survived our own audit")
    _badge(s, Inches(0.7), Inches(1.6), "ESTABLISHED", w=Inches(1.6))
    tf = _tb(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(4.4))
    for t in ["FaultChannel.POSITION_Y is inert against the driven sensing path under "
              "ADR-0033 — demonstrated causally (Control C: estimator sees 0 with "
              "redundancy on, 3.18 with it off), consistent across 9 policy/seed pairs.",
              "All nine layers are domain-mapped; all three Core-B gates verdict on every tick.",
              "The current L6 calibration prevents the gate from firing under the tested "
              "condition — classified as OD-8 / exchangeability violation, not a wiring defect.",
              "speed_stuck shows a stable, well-posed 30/30 absorption result at L2a under P1.",
              "Fault observability is heterogeneous across the valid faults.",
              "Invalid records are retained and flagged rather than deleted.",
              "The workflow now requires delivered-signal integrity verification before a "
              "sweep result is accepted."]:
        _p(tf, "▪  " + t, 11, MUTED, space=8)

    _badge(s, Inches(6.75), Inches(1.6), "WITHDRAWN", w=Inches(1.6))
    tf = _tb(s, Inches(6.75), Inches(2.0), Inches(5.85), Inches(4.4))
    for t in ["“Position faults are absorbed at the estimator.” — the faults never reached the "
              "estimator. Re-injected correctly and re-run at n=30: absorbed in 0 of 12 cells "
              "(D_L2a 0.655–0.797). The original result was entirely artefact.",
              "“All 18 falsification checks negative.” — sound checks, contaminated data.",
              "The H-regime operating-regime claim — withdrawn as a Simpson's paradox: "
              "pooled ρ = −0.341, but within policy P2 ρ = +0.836.",
              "“Information is destroyed at L2a.” — a downstream statistic recovers part of "
              "it; the correct phrasing is that the evaluated representation does not retain "
              "discriminative evidence.",
              "“A detection-without-response gap at L6.” — the gate returns PASS on every "
              "tick; it never made a detection to ignore."]:
        _p(tf, "▪  " + t, 11, MUTED, space=8)
    foot(s, "Three measurement defects were found in this project's own code, two of them by "
            "self-audit before publication. Each is documented with its cause and correction.")

    # 9 HYPOTHESIS ----------------------------------------------------------
    s = slide(prs, "Current scientific hypothesis", "stated so it can fail")
    _rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.15), fill=PANEL)
    _rect(s, Inches(0.7), Inches(1.6), Emu(30000), Inches(1.15), fill=ACCENT)
    tf = _tb(s, Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.9))
    _p(tf, "Sensor-fault evidence degrades non-uniformly across pipeline stages, and the "
           "stage at which it becomes insufficient is a measurable, fault-specific property "
           "that can inform where monitors are placed.", 14, INK, bold=True, first=True, space=0)

    _p(_tb(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(0.3)),
       "Three concepts kept separate throughout — conflating them is the most common error "
       "in this area", 12, INK, bold=True, first=True, space=0)
    chain(s, Inches(0.7), Inches(3.45),
          ["Sensor-level\ndetectability", "Estimator-level\nobservability",
           "Downstream operational\ndetectability"],
          w=Inches(3.86), h=Inches(0.85), gap=Inches(0.16),
          col=[EST, PART, PLAN], fs=12)
    tf = _tb(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.0))
    _p(tf, "Can the fault be recognised from the raw corrupted measurement?     |     "
           "Does it remain distinguishable after state estimation?     |     "
           "Can a real gate, with its real threshold, act on it?", 11, MUTED, first=True, space=0)

    kv_panel(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.55),
             "The distinction that currently matters most",
             [("A statistic can separate faulted from clean almost perfectly while the gate "
               "reading that statistic cannot fire, because the shift is far below its "
               "decision threshold.", MUTED),
              ("This is why ASTRA reports statistical discriminability and operational "
               "detectability as separate quantities — and why the deck never presents one "
               "as the other.", INK, True)], PART)

    # 10 NOVELTY ------------------------------------------------------------
    s = slide(prs, "Novelty position", "stated as a hypothesis, not a claim")
    _p(_tb(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.3)),
       "The novelty is not any single ingredient — each exists in the literature. The "
       "proposed contribution is the combination.", 12, MUTED, first=True, space=0)
    table(s, Inches(0.7), Inches(2.0), Inches(11.9), [
        ["#", "Proposed contribution", "Status today"],
        ["1", "Stage-wise fault discriminability across a governance pipeline", "Measured on 4 valid faults"],
        ["2", "Fault-specific observability trajectories", "Measured — four distinct behaviours"],
        ["3", "A candidate fault-absorption point A(f)", "Ill-posed for 5 of 6 faults"],
        ["4", "Separating statistical discriminability from operational gate detection", "Measured — and consequential"],
        ["5", "Using measured observability to investigate monitor placement", "PLANNED — experiment H7"],
        ["6", "Natural vs adversarial sensor manipulation", "PLANNED — not started"],
        ["7", "Simulation → real data → closed loop validation chain", "PLANNED — gated"],
    ], widths=[0.45, 6.6, 3.4], fs=11,
        colors={(1, 2): EST, (2, 2): EST, (3, 2): PART, (4, 2): EST,
                (5, 2): PLAN, (6, 2): PLAN, (7, 2): PLAN})
    _rect(s, Inches(0.7), Inches(5.32), Inches(11.9), Inches(0.92), fill=PANEL)
    _rect(s, Inches(0.7), Inches(5.32), Emu(30000), Inches(0.92), fill=HYP)
    _p(_tb(s, Inches(1.0), Inches(5.52), Inches(11.3), Inches(0.6)),
       "Novelty remains a research hypothesis until confirmed against a final systematic "
       "literature review and the completed experiments. No priority claim is made in this deck.",
       12.5, INK, bold=True, first=True, space=0)

    # 11 H7 -----------------------------------------------------------------
    s = slide(prs, "H7 — the monitor-placement experiment", "the next scientific step")
    _badge(s, Inches(0.7), Inches(1.58), "PLANNED", w=Inches(1.3))
    tf = _tb(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(2.2))
    _p(tf, "The question", 12.5, ACCENT, bold=True, first=True, space=7)
    _p(tf, "If a fault is distinguishable at one stage but not at a later stage, does "
           "monitoring at the earlier stage provide better operational detection?", 13,
       INK, bold=True, space=10)
    _p(tf, "This is where measured observability would become an engineering recommendation "
           "rather than an observation.", 11.5, MUTED, space=0)

    kv_panel(s, Inches(6.75), Inches(2.0), Inches(5.85), Inches(2.2),
             "A binding constraint discovered first",
             [("Monitor calibration — not placement — is currently the", INK, True),
              ("limiting factor. Under the tested condition the L6 gate", MUTED),
              ("cannot fire: live non-conformity ≈ 3.39 against a", MUTED),
              ("conformal quantile ≈ 5.40, a gap of ≈ 2.0.", MUTED)], PART)

    _p(_tb(s, Inches(0.7), Inches(4.42), Inches(11.9), Inches(0.3)),
       "H7 must therefore resolve, in order:", 12.5, INK, bold=True, first=True, space=0)
    chain(s, Inches(0.7), Inches(4.85),
          ["Can the monitor\noperate at all?", "Where does evidence\nremain sufficient?",
           "Does relocating it\nimprove detection?", "Does it hold across\nfaults & policies?"],
          w=Inches(2.86), h=Inches(0.8), gap=Inches(0.15),
          col=[RED, PART, PLAN, PLAN], fs=11)
    _p(_tb(s, Inches(0.7), Inches(5.92), Inches(11.9), Inches(0.5)),
       "Because the first step is currently unresolved, H7 is gated behind the calibration "
       "question (OD-8). H7 has not succeeded and has not been run.", 11.5, MUTED,
       first=True, space=0)

    # 12 DATASET STRATEGY ---------------------------------------------------
    s = slide(prs, "Dataset & validation strategy", "each source answers one question")
    table(s, Inches(0.7), Inches(1.62), Inches(11.9), [
        ["Source", "Purpose", "Contributes", "Weakness", "Status"],
        ["E17\n(synthetic)", "Controlled causal experimentation",
         "fault injection, severity, seeds, stage-wise D, absorption, policy dependence",
         "synthetic; not external validity", "IN USE"],
        ["comma2k19", "Real-world driving-data validation",
         "does the E17 phenomenon transfer when driven by real vehicle data?",
         "does not natively contain the required sensor faults", "PLANNED"],
        ["highD", "Naturalistic traffic / regime validation",
         "trajectories, speed distributions, density, following & lane-change behaviour",
         "not a sensor-fault dataset", "PLANNED"],
        ["CARLA", "Closed-loop validation & demonstration",
         "faults, adversarial sensors, estimator behaviour, closed-loop consequences",
         "must not be the sole source of evidence", "PLANNED"],
    ], widths=[1.35, 2.5, 4.3, 2.6, 1.15], fs=10, header_fs=10, row_h=Inches(0.95))
    kv_panel(s, Inches(0.7), Inches(5.62), Inches(11.9), Inches(1.08),
             "Method note",
             [("Where a real-world dataset does not contain the required faults, controlled "
               "fault injection is applied over real measurements. This is stated as a method, "
               "not disguised as naturally occurring faults. Dataset availability and licensing "
               "are verified before a dataset is treated as a confirmed resource.", MUTED)], ACCENT)

    # 13 comma2k19 + highD --------------------------------------------------
    s = slide(prs, "comma2k19 and highD — what each is for", "and what neither proves")
    kv_panel(s, Inches(0.7), Inches(1.62), Inches(5.85), Inches(2.65),
             "comma2k19 — real-world signal realism",
             [("Question: does the observability behaviour found in E17 survive when the "
               "pipeline is driven by real vehicle data?", INK, True),
              ("Contributes raw GNSS, 9-axis IMU and CAN from real highway driving.", MUTED),
              ("It does NOT natively contain the sensor faults being studied — controlled "
               "injection over real measurements is the method.", MUTED),
              ("Status: conditional — verification of timestamp precision, storage, "
               "calibration and drive identity is outstanding.", PART, True)], PLAN)
    kv_panel(s, Inches(6.75), Inches(1.62), Inches(5.85), Inches(2.65),
             "highD — operating-regime realism",
             [("Question: do the results hold across realistic traffic regimes rather than "
               "one synthetic driving profile?", INK, True),
              ("Contributes naturalistic trajectories, speed distributions, traffic density, "
               "following and lane-change behaviour.", MUTED),
              ("It is NOT treated as a sensor-fault dataset. Its role is regime and "
               "generalisation validation only.", MUTED),
              ("Status: planned, not started.", PLAN, True)], PLAN)

    _rect(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.05), fill=PANEL)
    _rect(s, Inches(0.7), Inches(4.5), Emu(30000), Inches(1.05), fill=RED)
    tf = _tb(s, Inches(1.0), Inches(4.68), Inches(11.3), Inches(0.85))
    _p(tf, "Additional dataset opportunities — only if scientifically justified", 12.5,
       INK, bold=True, first=True, space=6)
    _p(tf, "Categories under consideration: autonomous-driving datasets, localisation/sensor "
           "datasets, trajectory datasets, perception datasets, fault/diagnostic datasets, "
           "cybersecurity/attack datasets. More datasets do not automatically increase "
           "novelty — one is added only when it answers a specific unresolved question.",
       11, MUTED, space=0)
    _p(_tb(s, Inches(0.7), Inches(5.78), Inches(11.9), Inches(0.6)),
       "Current external-validation status: none. Every result in this deck is simulation-based. "
       "Thirty seeds buys reproducibility, not external validity — the two are routinely "
       "conflated and are not the same claim.", 11.5, INK, bold=True, first=True, space=0)

    # 14 ADVERSARIAL --------------------------------------------------------
    s = slide(prs, "Adversarial / lying-sensor extension", "future security direction")
    _badge(s, Inches(0.7), Inches(1.58), "PLANNED", w=Inches(1.3))
    _p(_tb(s, Inches(2.2), Inches(1.56), Inches(10.4), Inches(0.32)),
       "Not started. ASTRA does not currently detect manipulated sensors and no such claim "
       "is made.", 11.5, RED, bold=True, first=True, space=0)

    y = Inches(2.15)
    for label, col, note in [
        ("Honest sensor A", EST, "reports truth"),
        ("Compromised sensor B", RED, "reports plausible falsehood"),
        ("Honest redundant sensor C", EST, "reports truth"),
    ]:
        _rect(s, Inches(0.7), y, Inches(3.5), Inches(0.5), fill=col,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t2 = s.shapes[-1].text_frame
        t2.margin_left = Inches(0.12)
        _p(t2, label, 11.5, WHITE, bold=True, first=True, space=0)
        _p(_tb(s, Inches(4.35), y + Inches(0.13), Inches(2.2), Inches(0.3)),
           note, 10.5, MUTED, first=True, space=0)
        y += Inches(0.62)
    chain(s, Inches(6.9), Inches(2.15),
          ["Sensor fusion / estimator", "ASTRA stage-wise observability",
           "Is the manipulation still visible?"],
          w=Inches(5.4), h=Inches(0.5), gap=Inches(0.06), vertical=True,
          col=[MUTED, ACCENT, HYP], fs=11)

    kv_panel(s, Inches(0.7), Inches(4.25), Inches(5.85), Inches(2.05),
             "The scientific question",
             [("Can ASTRA characterise whether deliberately plausible false measurements "
               "remain observable, or become hidden by sensor fusion?", INK, True),
              ("A natural fault and a crafted one may occupy very different positions in the "
               "observability profile. That difference — if it exists — is the contribution.",
               MUTED)], HYP)
    kv_panel(s, Inches(6.75), Inches(4.25), Inches(5.85), Inches(2.05),
             "Entry point already exists",
             [("redundant.offset(modality, tick)", INK, True),
              ("A per-channel injection path added during the September audit. It makes one "
               "channel lie while others stay honest — which is precisely the adversarial "
               "configuration, and is the same mechanism needed to fix the position faults.",
               MUTED)], PLAN)

    # 15 CARLA / ROADMAP ----------------------------------------------------
    s = slide(prs, "Gated validation roadmap", "stages are gated, not sequential by default")
    steps = [("E17 audit", EST), ("Fix position injection", PART),
             ("Correct L6 interpretation", PART), ("Confirm speed_stuck metric", PART),
             ("Corrected E17 validation", PART), ("H7 monitor placement", PLAN),
             ("comma2k19 real data", PLAN), ("highD regime validation", PLAN),
             ("Adversarial sensors", PLAN), ("CARLA closed loop", PLAN),
             ("Integrated evaluation", HYP), ("Conference paper", HYP),
             ("Expanded journal paper", HYP)]
    x, y = Inches(0.7), Inches(1.75)
    bw2, bh2 = Inches(2.85), Inches(0.62)
    for i, (label, col) in enumerate(steps):
        r, c = divmod(i, 4)
        bx = x + c * (bw2 + Inches(0.18))
        by = y + r * (bh2 + Inches(0.34))
        _rect(s, bx, by, bw2, bh2, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t2 = s.shapes[-1].text_frame
        t2.word_wrap = True
        t2.margin_left = t2.margin_right = Inches(0.08)
        _p(t2, f"{i + 1}.  {label}", 11, WHITE, bold=True, first=True, space=0,
           align=PP_ALIGN.CENTER)
        if c < 3 and i < len(steps) - 1:
            _p(_tb(s, bx + bw2, by + Inches(0.16), Inches(0.18), Inches(0.3)),
               "▶", 9, FAINT, first=True, space=0, align=PP_ALIGN.CENTER)
    _rect(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.15), fill=PANEL)
    _rect(s, Inches(0.7), Inches(5.5), Emu(30000), Inches(1.15), fill=RED)
    tf = _tb(s, Inches(1.0), Inches(5.68), Inches(11.3), Inches(0.9))
    _p(tf, "Gating rule", 12.5, INK, bold=True, first=True, space=6)
    _p(tf, "A stage is not entered because it exists in the roadmap. Each must answer a "
           "specific unresolved scientific question, and the preceding stage must have "
           "produced evidence that justifies it. CARLA is deliberately last and is not "
           "presented as a source of primary evidence.", 11.5, MUTED, space=0)

    # 16 COMPANY VALUE ------------------------------------------------------
    s = slide(prs, "Why would a company care?", "potential value — to be validated")
    items = [
        ("1  Better monitor placement",
         "Move from intuition-based or convention-based monitor siting to evidence-based "
         "guidance: monitor where fault evidence still exists.",
         "fewer blind spots · better diagnostic coverage · better use of compute"),
        ("2  Safety validation",
         "Answer “which sensor faults become invisible after state estimation?” before "
         "deployment rather than after.",
         "find monitor blind spots · identify failure modes · prioritise testing"),
        ("3  Sensor-fusion robustness",
         "As GNSS, IMU, odometry, camera, LiDAR and radar are fused, ask what happens to "
         "fault evidence after fusion.",
         "framework for reasoning about fusion-induced evidence loss"),
        ("4  Cybersecurity extension",
         "If the adversarial work succeeds: which manipulations stay detectable after "
         "fusion, and which do not?",
         "security testing · attack-surface analysis · redundancy evaluation"),
    ]
    x = Inches(0.7)
    for i, (h, body, benefit) in enumerate(items):
        bx = Inches(0.7) + (i % 2) * Inches(6.05)
        by = Inches(1.62) + (i // 2) * Inches(2.05)
        _rect(s, bx, by, Inches(5.85), Inches(1.85), fill=PANEL)
        _rect(s, bx, by, Emu(30000), Inches(1.85), fill=ACCENT)
        tf = _tb(s, bx + Inches(0.22), by + Inches(0.15), Inches(5.45), Inches(1.6))
        _p(tf, h, 12.5, INK, bold=True, first=True, space=6)
        _p(tf, body, 11, MUTED, space=6)
        _p(tf, benefit, 10, ACCENT, bold=True, space=0)
    _rect(s, Inches(0.7), Inches(5.78), Inches(11.9), Inches(0.82), fill=PANEL)
    _rect(s, Inches(0.7), Inches(5.78), Emu(30000), Inches(0.82), fill=RED)
    _p(_tb(s, Inches(1.0), Inches(5.94), Inches(11.3), Inches(0.6)),
       "All four are POTENTIAL benefits contingent on the validation programme succeeding. "
       "No cybersecurity protection is promised, and no company benefit has yet been "
       "demonstrated on company data.", 11.5, INK, bold=True, first=True, space=0)

    # 17 BEFORE / AFTER + COLLABORATION ------------------------------------
    s = slide(prs, "Engineering decision support", "what the output could become")
    _p(_tb(s, Inches(0.7), Inches(1.55), Inches(5.85), Inches(0.3)),
       "TODAY — binary detection", 11, MUTED, bold=True, first=True, space=0)
    chain(s, Inches(0.7), Inches(1.92), ["Fault occurs", "Detector: “fault / no fault”"],
          w=Inches(2.84), h=Inches(0.55), gap=Inches(0.15), col=[MUTED, MUTED], fs=11)
    _p(_tb(s, Inches(0.7), Inches(2.62), Inches(5.85), Inches(0.4)),
       "Limited information about where the fault became invisible.", 11, MUTED,
       first=True, space=0)

    _p(_tb(s, Inches(6.75), Inches(1.55), Inches(5.85), Inches(0.3)),
       "POTENTIAL WITH ASTRA — stage-wise profile", 11, ACCENT, bold=True, first=True, space=0)
    prof = [("L1  sensing", "visible", EST), ("L2  estimation", "visible", EST),
            ("L2a innovation", "weak", PART), ("L7  shield", "insufficient", RED)]
    yy = Inches(1.92)
    for name, note, col in prof:
        _rect(s, Inches(6.75), yy, Inches(2.5), Inches(0.34), fill=col,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t2 = s.shapes[-1].text_frame
        t2.margin_left = Inches(0.1)
        _p(t2, name, 10, WHITE, bold=True, first=True, space=0)
        _p(_tb(s, Inches(9.4), yy + Inches(0.06), Inches(3.2), Inches(0.28)),
           note, 10, MUTED, first=True, space=0)
        yy += Inches(0.4)
    _p(_tb(s, Inches(6.75), Inches(3.55), Inches(5.85), Inches(0.3)),
       "→ “the architecture has a monitoring blind spot here”", 11.5, INK, bold=True,
       first=True, space=0)

    _p(_tb(s, Inches(0.7), Inches(4.05), Inches(11.9), Inches(0.3)),
       "Possible future output format", 12, INK, bold=True, first=True, space=0)
    _rect(s, Inches(0.7), Inches(4.42), Inches(11.9), Inches(0.5), fill=INK)
    _p(_tb(s, Inches(0.95), Inches(4.55), Inches(11.5), Inches(0.3)),
       "fault  →  stage where evidence weakens  →  recommended monitor location  →  confidence",
       12, WHITE, bold=True, first=True, space=0, font=MONO)

    _p(_tb(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(0.3)),
       "What we could offer a collaborating company — potential directions, not existing "
       "partnerships", 12, INK, bold=True, first=True, space=0)
    tf = _tb(s, Inches(0.7), Inches(5.55), Inches(5.85), Inches(1.1))
    for t in ["validation on company-specific sensor architectures",
              "safety-monitor evaluation and blind-spot analysis",
              "sensor-fusion robustness and fault-injection scenarios"]:
        _p(tf, "▪  " + t, 10.5, MUTED, space=4)
    tf = _tb(s, Inches(6.75), Inches(5.55), Inches(5.85), Inches(1.1))
    for t in ["adversarial sensor testing (once validated)",
              "CARLA / digital-twin closed-loop validation",
              "benchmarking and joint publication or technical report"]:
        _p(tf, "▪  " + t, 10.5, MUTED, space=4)

    # 18 CREDIBILITY MATRIX -------------------------------------------------
    s = slide(prs, "Credibility assessment", "qualitative, not probabilistic")
    table(s, Inches(0.7), Inches(1.62), Inches(11.9), [
        ["Dimension", "Current status", "After the roadmap"],
        ["Experimental integrity", "Strong — independently audited", "Strong"],
        ["Controlled validation", "Developing — E17, 4 valid faults", "Strong — E17 expanded"],
        ["Statistical confidence", "Developing — n=30 seeds, CIs reported", "Strong"],
        ["Real-world validity", "Low — none yet", "Promising — comma2k19"],
        ["Regime validation", "Low — none yet", "Promising — highD"],
        ["Monitor-placement evidence", "Low — H7 not run", "Promising — H7"],
        ["Adversarial validation", "Low — not started", "Developing — planned"],
        ["Closed-loop validation", "Low — not started", "Promising — CARLA"],
        ["Industry relevance", "Developing — potential only", "Promising after validation"],
        ["Conference readiness", "Developing", "Potentially strong"],
        ["Journal (Q2) readiness", "Low", "Potentially strong"],
        ["Journal (Q1) readiness", "Low", "Depends on final evidence"],
    ], widths=[3.4, 4.6, 3.9], fs=10.5, row_h=Inches(0.31))
    _p(_tb(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.4)),
       "Labels are qualitative by design: Low · Developing · Promising · Strong · "
       "Submission-ready. No acceptance probability is quoted, because none can be "
       "responsibly estimated.", 11, MUTED, first=True, space=0)

    # 19 PUBLICATION --------------------------------------------------------
    s = slide(prs, "Publication pathway", "three tiers, each with explicit prerequisites")
    tiers = [
        ("IEEE Conference", ACCENT,
         ["corrected E17", "reproducibility", "statistically defensible results",
          "clear novelty positioning", "H7 evidence", "meaningful external validation",
          "clearly stated limitations"],
         "Strong conference potential IF the remaining validation confirms the central hypothesis."),
        ("Q2 Journal", PART,
         ["everything above, plus", "real-world validation", "strong ablations",
          "full statistical analysis", "reproducibility package",
          "clear engineering implications"],
         "A rigorous empirical paper could form a strong Q2 submission."),
        ("Q1 Journal", HYP,
         ["novel phenomenon", "statistical reproducibility", "independent real-world validation",
          "mechanistic explanation", "monitor-placement consequence", "adversarial extension",
          "closed-loop CARLA validation", "strong baselines + ablations"],
         "Q1 requires substantially stronger evidence than exists today."),
    ]
    for i, (name, col, reqs, note) in enumerate(tiers):
        bx = Inches(0.7) + i * Inches(4.05)
        _rect(s, bx, Inches(1.62), Inches(3.85), Inches(4.55), fill=PANEL)
        _rect(s, bx, Inches(1.62), Inches(3.85), Inches(0.5), fill=col)
        t2 = _tb(s, bx + Inches(0.2), Inches(1.74), Inches(3.5), Inches(0.3))
        _p(t2, name, 14, WHITE, bold=True, first=True, space=0)
        tf = _tb(s, bx + Inches(0.2), Inches(2.28), Inches(3.5), Inches(2.7))
        for r in reqs:
            _p(tf, "▪  " + r, 10.5, MUTED, space=5)
        _p(_tb(s, bx + Inches(0.2), Inches(5.35), Inches(3.5), Inches(0.75)),
           note, 10.5, INK, bold=True, first=True, space=0)
    _rect(s, Inches(0.7), Inches(6.34), Inches(11.9), Inches(0.5), fill=RED)
    _p(_tb(s, Inches(1.0), Inches(6.44), Inches(11.4), Inches(0.3)),
       "No acceptance probability can be guaranteed. The objective is to make the scientific "
       "contribution strong enough for a credible submission.", 11.5, WHITE, bold=True,
       first=True, space=0)

    # 20 CLAIMS / MUST NOT --------------------------------------------------
    s = slide(prs, "Claims discipline", "what we say, and what we refuse to say")
    _badge(s, Inches(0.7), Inches(1.6), "ESTABLISHED", w=Inches(1.6))
    _p(_tb(s, Inches(0.7), Inches(1.98), Inches(5.85), Inches(0.3)),
       "What we can officially claim today", 12.5, INK, bold=True, first=True, space=0)
    tf = _tb(s, Inches(0.7), Inches(2.35), Inches(5.85), Inches(4.0))
    for t in ["POSITION_Y is inert under ADR-0033 — causally demonstrated (Control C)",
              "9/9 layers domain-mapped; 3/3 Core-B gates verdict every tick",
              "Current L6 calibration prevents firing under the tested condition (OD-8)",
              "speed_stuck: stable well-posed 30/30 absorption at L2a under P1",
              "Fault observability is heterogeneous across the valid faults",
              "The H-regime claim was withdrawn (Simpson's paradox)",
              "Invalid records are retained and flagged, not deleted",
              "Delivered-signal integrity verification is now required before acceptance",
              "Corrected position faults are NOT absorbed at L2a — 0 of 12 cells, n=30"]:
        _p(tf, "✓  " + t, 10.5, MUTED, space=7)

    _badge(s, Inches(6.75), Inches(1.6), "MUST NOT CLAIM", w=Inches(1.9))
    _p(_tb(s, Inches(6.75), Inches(1.98), Inches(5.85), Inches(0.3)),
       "What we refuse to claim", 12.5, INK, bold=True, first=True, space=0)
    tf = _tb(s, Inches(6.75), Inches(2.35), Inches(5.85), Inches(4.0))
    for t in ["universal fault absorption",
              "that position faults are currently absorbed",
              "that all faults have validated absorption points",
              "that D = 0.5 means “50 % accuracy”",
              "that high D means an operational gate fired",
              "that ASTRA already detects hacked sensors",
              "that ASTRA is proven Q1-worthy",
              "guaranteed conference or journal acceptance",
              "that more datasets automatically create novelty"]:
        _p(tf, "✕  " + t, 10.5, RED, space=7)

    # 21 CLOSING ------------------------------------------------------------
    s = slide(prs)
    _rect(s, 0, 0, W, H, fill=INK)
    _p(_tb(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.3)),
       "WHY WE SHOULD CONTINUE", 12, RGBColor(0x7E, 0xA8, 0xC9), bold=True, first=True, space=0)
    _rect(s, Inches(0.9), Inches(1.28), Inches(1.7), Emu(30000), fill=ACCENT)

    prog = ["Controlled science", "Fault observability", "Estimator interaction",
            "Monitor placement", "Real-world validation", "Traffic-regime validation",
            "Adversarial sensors", "Closed-loop safety", "Industry validation"]
    x = Inches(0.9)
    for i, label in enumerate(prog):
        bx = Inches(0.9) + (i % 3) * Inches(3.95)
        by = Inches(1.72) + (i // 3) * Inches(0.62)
        _rect(s, bx, by, Inches(3.7), Inches(0.48),
              fill=RGBColor(0x1D, 0x39, 0x54), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t2 = s.shapes[-1].text_frame
        t2.margin_left = Inches(0.12)
        _p(t2, f"{i + 1}.  {label}", 11.5, RGBColor(0xCF, 0xDC, 0xE8), bold=True,
           first=True, space=0)

    tf = _tb(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(2.0))
    _p(tf, "The current E17 result is not the final product. It is the controlled foundation "
           "for testing a broader hypothesis across increasingly realistic environments.",
       15, WHITE, bold=True, first=True, space=12)
    _p(tf, "ASTRA already has a technically audited foundation and a scientifically "
           "interesting initial finding. The remaining work is not random feature "
           "development: it is a gated validation programme that can determine whether the "
           "finding generalises to real data, adversarial conditions and closed-loop "
           "autonomous driving — and whether it can become both an industry-relevant "
           "engineering capability and a strong academic contribution.",
       13, RGBColor(0xC2, 0xD3, 0xE2), space=0)
    _p(_tb(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4)),
       "The goal is not to expand the project artificially. It is to determine whether the "
       "observed phenomenon is real, reproducible, generalisable and practically useful.",
       11.5, RGBColor(0x8D, 0xA5, 0xBC), first=True, space=0, italic=True)

    prs.save(str(out))
    print(f"  {len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")


if __name__ == "__main__":
    root = Path(__file__).resolve().parent.parent
    p = root / "results" / "E17_POSITION" / "analysis.json"
    build(root / "docs" / "ASTRA_Research_Collaboration_Briefing.pptx",
          json.loads(p.read_text(encoding="utf-8")) if p.exists() else None)

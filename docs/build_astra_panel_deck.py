"""ASTRA Major Project Panel Review deck - 18 slides, dense text.

Every technical claim traceable to research/, experiments/phase5_od8_h7/, or results/.
Nothing invented. [NEEDS INPUT] marks fields that require team input.

Run with the isolated docs venv, NOT the measurement venv.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

W, H = Inches(13.333), Inches(7.5)

# Charcoal Minimal + Teal Trust - restrained, panel-appropriate
INK = RGBColor(0x1A, 0x22, 0x2E)
BODY = RGBColor(0x38, 0x45, 0x53)
MUTE = RGBColor(0x5F, 0x6E, 0x7C)
FAINT = RGBColor(0x93, 0xA0, 0xAB)
LINE = RGBColor(0xD5, 0xDC, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xF3, 0xF6, 0xF8)
TINT2 = RGBColor(0xE7, 0xED, 0xF1)

ACCENT = RGBColor(0x02, 0x80, 0x90)         # Teal Trust
ACCENT_DK = RGBColor(0x01, 0x5F, 0x6B)
HEAD_BG = RGBColor(0x14, 0x28, 0x32)

OK = RGBColor(0x1E, 0x6B, 0x52)
WARN = RGBColor(0xA9, 0x76, 0x0B)
BAD = RGBColor(0x9B, 0x33, 0x25)
INPT = RGBColor(0x6B, 0x2E, 0x86)           # for [NEEDS INPUT] chips

CHIP = {"DONE": OK, "IN PROGRESS": WARN, "NEXT": ACCENT, "PLANNED": MUTE,
        "BLOCKED": BAD, "WITHDRAWN": BAD, "NEEDS INPUT": INPT}

HEAD = "Cambria"
TEXT = "Calibri"


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    t = s.shapes.add_textbox(x, y, w, h).text_frame
    t.word_wrap = True
    t.vertical_anchor = anchor
    t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
    return t


def p(tf, text, size, color=BODY, bold=False, first=False, space=5, font=TEXT,
      align=PP_ALIGN.LEFT, italic=False):
    par = tf.paragraphs[0] if first else tf.add_paragraph()
    par.alignment = align
    par.space_after = Pt(space)
    r = par.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return par


def rect(s, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE, line=None):
    o = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        o.fill.background()
    else:
        o.fill.solid()
        o.fill.fore_color.rgb = fill
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = line
        o.line.width = Pt(0.75)
    o.shadow.inherit = False
    return o


def chip(s, x, y, label, w=None):
    col = CHIP.get(label, MUTE)
    w = w or Inches(0.16 + 0.083 * len(label))
    o = rect(s, x, y, w, Inches(0.24), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    t = o.text_frame
    t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
    t.word_wrap = False
    p(t, label, 8.5, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
    return w


def slide(prs, title=None, kicker=None, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = HEAD_BG if dark else WHITE
    if title:
        if kicker:
            p(tb(s, Inches(0.6), Inches(0.38), Inches(12.1), Inches(0.24)),
              kicker.upper(), 9.5, ACCENT, bold=True, first=True, space=0)
        p(tb(s, Inches(0.6), Inches(0.64), Inches(12.1), Inches(0.6)),
          title, 26, WHITE if dark else INK, bold=True, first=True, space=0, font=HEAD)
    return s


def table(s, x, y, w, rows, widths=None, fs=10, hfs=10, rh=Inches(0.28),
          hh=Inches(0.30), colors=None):
    t = s.shapes.add_table(len(rows), len(rows[0]), x, y, w, hh + rh * (len(rows) - 1)).table
    t.rows[0].height = hh
    for i in range(1, len(rows)):
        t.rows[i].height = rh
    if widths:
        tot = sum(widths)
        for j, fr in enumerate(widths):
            t.columns[j].width = Emu(int(w * fr / tot))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.margin_left = c.margin_right = Inches(0.07)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = HEAD_BG if i == 0 else (WHITE if i % 2 else TINT)
            tf = c.text_frame
            tf.word_wrap = True
            par = tf.paragraphs[0]
            par.space_after = Pt(0)
            r = par.add_run()
            r.text = str(val)
            r.font.size = Pt(hfs if i == 0 else fs)
            r.font.bold = i == 0
            r.font.name = TEXT
            if i == 0:
                r.font.color.rgb = WHITE
            elif colors and (i, j) in colors:
                r.font.color.rgb = colors[(i, j)]
                r.font.bold = True
            else:
                r.font.color.rgb = INK
    return t


def para(tf, items, size=10.5, color=BODY, space=3, prefix="•  "):
    for it in items:
        p(tf, prefix + it, size, color, space=space)


def heading(tf, text, size=12, color=ACCENT_DK, space=4, first=False):
    p(tf, text, size, color, bold=True, first=first, space=space, font=HEAD)


def foot(s, text):
    p(tb(s, Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.28)),
      text, 8.5, FAINT, first=True, space=0, italic=True)


def input_tag(s, x, y):
    chip(s, x, y, "NEEDS INPUT", w=Inches(1.15))


# ------------------------------------------------------------------ deck ----
def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 TITLE ---------------------------------------------------------------
    s = slide(prs, dark=True)
    rect(s, 0, 0, W, H, fill=HEAD_BG)
    rect(s, Inches(0.6), Inches(1.35), Inches(0.14), Inches(0.62), fill=ACCENT)
    p(tb(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.4)),
      "MAJOR PROJECT PANEL REVIEW  ·  DEPARTMENT OF COMPUTER SCIENCE & ENGINEERING",
      11, RGBColor(0x7F, 0xC5, 0xD1), bold=True, first=True, space=0)
    tf = tb(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(2.6))
    p(tf, "ASTRA", 60, WHITE, bold=True, first=True, space=2, font=HEAD)
    p(tf, "Autonomous Safety-Trust Runtime Architecture", 22,
      RGBColor(0xC6, 0xDD, 0xE3), space=6, font=HEAD)
    p(tf, "A layered runtime governance stack for learned controllers, and an "
          "experimental study of when its safety monitors can and cannot be trusted",
      13, RGBColor(0x94, 0xB8, 0xC3), italic=True, space=0)

    # team block
    tf = tb(s, Inches(0.9), Inches(5.05), Inches(6.4), Inches(1.9))
    p(tf, "PROJECT TEAM", 10, ACCENT, bold=True, first=True, space=6)
    p(tf, "Sushanth C.", 13, WHITE, bold=True, space=1)
    p(tf, "USN: [NEEDS INPUT]", 10.5, RGBColor(0x9C, 0xB6, 0xC0), space=6)
    p(tf, "Tarun Gowda V", 13, WHITE, bold=True, space=1)
    p(tf, "USN: [NEEDS INPUT]", 10.5, RGBColor(0x9C, 0xB6, 0xC0), space=6)
    p(tf, "T Tilak Reddy", 13, WHITE, bold=True, space=1)
    p(tf, "USN: [NEEDS INPUT]", 10.5, RGBColor(0x9C, 0xB6, 0xC0), space=0)

    tf = tb(s, Inches(7.4), Inches(5.05), Inches(5.1), Inches(1.9))
    p(tf, "GUIDE", 10, ACCENT, bold=True, first=True, space=6)
    p(tf, "Dr. Chaitra R.", 13, WHITE, bold=True, space=1)
    p(tf, "Department of Computer Science & Engineering", 10.5,
      RGBColor(0x9C, 0xB6, 0xC0), space=14)
    p(tf, "INSTITUTION", 10, ACCENT, bold=True, space=6)
    p(tf, "BMS College of Engineering, Bengaluru", 13, WHITE, bold=True, space=1)
    p(tf, "Academic Year 2025-2026", 10.5, RGBColor(0x9C, 0xB6, 0xC0), space=0)

    # 2 PROJECT OVERVIEW ----------------------------------------------------
    s = slide(prs, "Project overview", "at a glance")
    tf = tb(s, Inches(0.6), Inches(1.4), Inches(8.0), Inches(5.4))
    heading(tf, "What ASTRA is", first=True)
    p(tf, "A nine-layer runtime safety architecture for a vehicle driven by a learned "
          "(reinforcement-learning) controller, together with an experimental testbed "
          "that injects controlled sensor faults and measures how each layer detects — or "
          "fails to detect — those faults.", 11.5, BODY, space=10)

    heading(tf, "What problem it addresses")
    p(tf, "Autonomous cyber-physical systems increasingly rely on learned controllers whose "
          "internal behaviour cannot be formally verified. Runtime safety monitors are placed "
          "around them to catch sensor faults and misbehaviour. Existing architectures assume "
          "these monitors work reliably. Our experiments show that this assumption fails in "
          "specific, measurable ways — including a monitor being silent during the very "
          "hazard it is designed to catch. ASTRA characterises these failures and proposes an "
          "architecture that assesses the trustworthiness of its own monitors at runtime.",
      11.5, BODY, space=10)

    heading(tf, "Intended users")
    para(tf, [
        "Autonomous-vehicle and ADAS engineers designing runtime safety supervisors",
        "Safety analysts performing SOTIF (ISO 21448) coverage assessments",
        "Researchers in runtime verification and conformal prediction for control systems",
    ], size=11, space=3)

    # side panel
    rect(s, Inches(9.0), Inches(1.4), Inches(3.7), Inches(5.4), fill=TINT)
    tf = tb(s, Inches(9.2), Inches(1.55), Inches(3.4), Inches(5.1))
    p(tf, "PROJECT AT A GLANCE", 10, ACCENT, bold=True, first=True, space=8)
    p(tf, "6", 30, INK, bold=True, space=0, font=HEAD)
    p(tf, "experiment phases completed", 10, MUTE, space=10)
    p(tf, "~4,300", 30, INK, bold=True, space=0, font=HEAD)
    p(tf, "closed-loop simulation runs", 10, MUTE, space=10)
    p(tf, "9", 30, INK, bold=True, space=0, font=HEAD)
    p(tf, "architectural layers, verified", 10, MUTE, space=10)
    p(tf, "6", 30, INK, bold=True, space=0, font=HEAD)
    p(tf, "measurement defects self-caught", 10, MUTE, space=0)
    foot(s, "All results are from a synthetic driving plant. External real-world validation is not "
            "yet in scope for the current stage.")

    # 3 PROBLEM STATEMENT --------------------------------------------------
    s = slide(prs, "Problem statement", "the real-world problem")
    tf = tb(s, Inches(0.6), Inches(1.4), Inches(7.6), Inches(5.5))
    heading(tf, "The problem, precisely stated", first=True)
    p(tf, "A learned controller in a vehicle depends on sensor input. Sensors fail — bias, "
          "drift, dropout, noise, stuck values. Around the learned controller a safety "
          "supervisor runs runtime monitors intended to detect these failures and switch to "
          "a safe fallback before harm occurs.", 11.5, BODY, space=8)
    p(tf, "The safety of the whole system therefore depends on the monitors being reliable. "
          "Every published safety architecture we surveyed assumes this without measurement.",
      11.5, BODY, bold=True, space=10)

    heading(tf, "Why this matters")
    para(tf, [
        "Safety standards (ISO 26262 functional safety, ISO 21448 SOTIF) require coverage "
        "arguments for how faults are detected — those arguments rest on monitor reliability.",
        "A monitor that is silent during a hazard is worse than one that misses it: silence "
        "reads as evidence of health, so the safe fallback is never engaged.",
        "There is no standard method to check whether a runtime monitor can, in principle, "
        "make a reliable decision before it is deployed.",
    ], size=11)

    heading(tf, "Why a technological solution is required")
    p(tf, "Analytical verification of learned controllers is not feasible at scale. The "
          "pragmatic response is runtime monitoring, but runtime monitoring is currently "
          "designed and deployed without a discipline for assessing whether the monitor "
          "itself is trustworthy under the operating conditions it will actually see.",
      11.5, BODY, space=0)

    rect(s, Inches(8.5), Inches(1.4), Inches(4.2), Inches(5.5), fill=TINT2)
    tf = tb(s, Inches(8.7), Inches(1.55), Inches(3.9), Inches(5.2))
    p(tf, "WHO EXPERIENCES IT", 10, ACCENT, bold=True, first=True, space=6)
    para(tf, [
        "Vehicle manufacturers integrating learned components in perception, prediction and "
        "control",
        "ADAS Tier-1 suppliers whose safety cases include runtime monitors",
        "Certification bodies asked to sign off on learned-controller safety arguments",
        "Downstream: passengers and vulnerable road users, whose safety depends on those "
        "monitors actually working",
    ], size=10, space=6)

    # 4 MOTIVATION ---------------------------------------------------------
    s = slide(prs, "Motivation", "why we selected this problem")
    tf = tb(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.5))
    heading(tf, "Real-world relevance", first=True)
    p(tf, "Learned controllers are already deployed at scale in production ADAS. Public data "
          "on their failure modes is limited, and the safety cases behind them lean heavily "
          "on runtime supervisors. If those supervisors have systematic blind spots, the "
          "assurance argument does not hold.",
      11, BODY, space=8)

    heading(tf, "Existing pain points")
    para(tf, [
        "The Simplex architecture (Sha, 2001) is the field-standard runtime-assurance pattern "
        "but predates learned controllers and assumes the decision logic works.",
        "Conformal prediction gives elegant statistical guarantees under exchangeability, but "
        "no established procedure checks whether that assumption holds in a live system.",
        "Fault-detection literature evaluates a monitor at one location under one metric — "
        "typically AUC — and does not measure whether operational firing follows.",
    ], size=10.5, space=5)

    heading(tf, "What motivated this system")
    p(tf, "We wanted a testbed where a fault could be traced through every layer of a "
          "governance stack, one variable at a time, with the discipline that any monitoring "
          "claim must survive delivered-signal verification and a pre-registered evaluation "
          "criterion. The testbed then produced the finding: monitors have systematic blind "
          "spots that no calibration alone repairs.",
      11, BODY, space=0)

    rect(s, Inches(7.0), Inches(1.4), Inches(5.7), Inches(5.5), fill=TINT)
    tf = tb(s, Inches(7.2), Inches(1.55), Inches(5.4), Inches(5.2))
    p(tf, "WHY EXISTING APPROACHES ARE INSUFFICIENT", 10, ACCENT, bold=True,
      first=True, space=8)
    rows = [
        ("Simplex (2001)", "prescribes fallback controller and switching logic; does not "
                           "prescribe how the decision logic itself is judged reliable"),
        ("Innovation-based FDI", "detects at one stage; not routinely evaluated against a "
                                 "learned controller's downstream behaviour"),
        ("Conformal prediction", "provides coverage under exchangeability; live systems "
                                 "routinely violate exchangeability without warning"),
        ("Digital-twin / model-based", "monitors residuals; assumes model and system stay "
                                       "aligned, an assumption that itself needs checking"),
        ("Redundant sensor fusion", "voting outvotes a liar and discards the evidence; "
                                    "fault correction becomes fault concealment"),
    ]
    for ttl, txt in rows:
        p(tf, ttl, 10.5, INK, bold=True, space=2)
        p(tf, txt, 9.5, BODY, space=8)

    # 5 OBJECTIVES ---------------------------------------------------------
    s = slide(prs, "Objectives", "what we set out to do")
    tf = tb(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.3))
    p(tf, "Six objectives, each traceable to a specific experimental deliverable in the "
          "project record.", 11, MUTE, first=True, space=0)
    objs = [
        ("1", "Design a layered runtime governance architecture", ACCENT,
         "Build a nine-layer stack around a learned proposer (Core-A) with a one-way trust "
         "boundary to a safety domain (Core-B). Verify every layer is domain-mapped and every "
         "safety gate returns a verdict on every control tick."),
        ("2", "Build a controlled fault-injection testbed", ACCENT,
         "Six fault classes across three sensor channels, delivered-signal verification, "
         "pre-registered evaluation windows and thresholds. Testbed used to run 4,300 "
         "closed-loop simulation runs."),
        ("3", "Characterise stage-wise fault observability", ACCENT,
         "For each fault and each layer, measure a stage-wise discriminability score. Produce "
         "an observability profile that shows where evidence of the fault survives and where "
         "it does not."),
        ("4", "Establish calibration and detection behaviour of a conformal safety monitor",
         WARN,
         "Systematically calibrate the statistical gate under three schemes; measure per-run "
         "false-alarm and detection behaviour; report the mechanism when it fails."),
        ("5", "Separate statistical discriminability from operational detectability",
         WARN,
         "Show, with data, whether a stage where the fault is statistically distinguishable "
         "also permits an operational monitor to fire reliably."),
        ("6", "Propose an architectural response — a self-trust layer",
         MUTE,
         "Design ASTRA 2.0, a monitoring architecture that assesses whether its own monitors "
         "can be trusted at runtime. Design published; implementation is future work."),
    ]
    y = Inches(1.85)
    for n, title, col, body in objs:
        rect(s, Inches(0.6), y, Inches(12.1), Inches(0.82), fill=TINT if int(n) % 2 else TINT2)
        rect(s, Inches(0.6), y, Inches(0.65), Inches(0.82), fill=col)
        p(tb(s, Inches(0.6), y + Inches(0.22), Inches(0.65), Inches(0.4)),
          n, 22, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER, font=HEAD)
        tf = tb(s, Inches(1.4), y + Inches(0.08), Inches(11.15), Inches(0.72))
        p(tf, title, 12, INK, bold=True, first=True, space=3, font=HEAD)
        p(tf, body, 10, BODY, space=0)
        y += Inches(0.87)
    foot(s, "Objectives 1-5 are complete or substantially complete; objective 6 is designed with "
            "a validation plan pending implementation.")

    # 6 EXISTING SYSTEMS + GAP --------------------------------------------
    s = slide(prs, "Existing systems and identified gap", "literature survey")
    table(s, Inches(0.6), Inches(1.4), Inches(12.1), [
        ["Approach", "Core idea", "Strength", "Limitation for our setting"],
        ["Simplex (Sha 2001)", "Complex controller + safety controller + switching logic",
         "Formalised trust boundary; widely adopted",
         "Assumes decision logic is reliable; no self-assessment"],
        ["Analytical redundancy / FDI\n(Willsky 1976; Basseville & Nikiforov)",
         "Detect faults from residuals or innovation statistics",
         "Well-studied, model-based, provable under linear assumptions",
         "Requires accurate model; typically single detector at one stage"],
        ["Conformal prediction / ICP\n(Vovk, Shafer, Gammerman)",
         "Non-parametric anomaly detection with coverage guarantees",
         "Distribution-free under exchangeability",
         "Exchangeability assumption routinely violated in live systems; no built-in check"],
        ["Digital-twin monitoring\n(model-based diagnosis)",
         "Compare live system output with a model of expected behaviour",
         "Interpretable; supports what-if analysis",
         "Twin fidelity itself must be maintained; drift is silent"],
        ["Redundant sensor fusion\n(TMR, median voting)",
         "Multiple sensors, majority vote rejects the outlier",
         "Handles common single-channel faults",
         "A well-crafted adversarial signal can hide inside the voting margin"],
        ["Fault-injection studies\n(various)",
         "Inject known faults, measure detector response",
         "Empirical, comparable across systems",
         "Typically evaluated at one stage; run-level statistics rarely reported"],
    ], widths=[1.9, 3.2, 3.2, 3.8], fs=9.5, hfs=9.5, rh=Inches(0.62))

    rect(s, Inches(0.6), Inches(6.02), Inches(12.1), Inches(0.75), fill=RGBColor(0xEE, 0xE0, 0xE8))
    rect(s, Inches(0.6), Inches(6.02), Inches(0.14), Inches(0.75), fill=INPT)
    tf = tb(s, Inches(0.9), Inches(6.10), Inches(11.7), Inches(0.65))
    p(tf, "IDENTIFIED GAP", 10, INPT, bold=True, first=True, space=3)
    p(tf, "Existing runtime safety architectures assume their monitors work. There is no "
          "established, reproducible procedure that measures whether a runtime monitor can be "
          "trusted under the conditions it will actually see — before deployment or during it. "
          "ASTRA addresses this gap directly and demonstrates its necessity empirically.",
      10.5, INK, space=0)
    foot(s, "References available in accompanying literature-review document; systematic review "
            "recorded as an in-progress deliverable.")

    # 7 PROPOSED SYSTEM ----------------------------------------------------
    s = slide(prs, "Proposed system", "how ASTRA addresses the gap")
    tf = tb(s, Inches(0.6), Inches(1.4), Inches(6.2), Inches(5.5))
    heading(tf, "The proposal in one sentence", first=True)
    p(tf, "A layered governance stack around a learned controller, plus a self-assessment "
          "mechanism that measures whether each monitor can currently be trusted, and routes "
          "decisions accordingly.", 12, INK, bold=True, space=10)

    heading(tf, "How it addresses the gap")
    para(tf, [
        "Separates two questions that existing systems collapse: is there evidence of a fault, "
        "and can a decision be made on that evidence?",
        "Instruments every layer — sensing, estimation, trust, proposer, twin, statistical gate, "
        "deterministic shield, fail-safe, arbitration — so fault evidence is measured, not assumed.",
        "Adds a monitorability metric that predicts whether a stage's decision statistic has the "
        "dynamic range and stability to support a reliable threshold, before deployment.",
        "Adds a phase-aware detection rule that distinguishes fault activity from post-fault "
        "recovery, so that a monitor's silence during the fault is not read as safety.",
        "Preserves the one-way trust boundary between the learned domain (Core-A) and the "
        "safety domain (Core-B) throughout.",
    ], size=10.5, space=4)

    heading(tf, "What makes the approach useful")
    p(tf, "The self-assessment mechanism is what makes an operator able to say, at runtime, "
          "how much to trust the current output of a specific monitor. That is a capability "
          "no existing architecture in this space provides.", 10.5, BODY, space=0)

    # visual: two-column diagram
    rect(s, Inches(7.0), Inches(1.4), Inches(5.7), Inches(5.5), fill=TINT)
    tf = tb(s, Inches(7.2), Inches(1.55), Inches(5.4), Inches(0.3))
    p(tf, "CONCEPTUAL VIEW", 10, ACCENT, bold=True, first=True, space=0)

    # Stack column
    y0 = Inches(1.95)
    for i, (code, name, col) in enumerate([
        ("L1", "sensing", ACCENT_DK), ("L2", "UKF estimator", ACCENT_DK),
        ("L3", "trust", ACCENT_DK), ("L4", "learned proposer  (CORE-A)", MUTE),
        ("L5", "PINN twin", ACCENT), ("L6", "conformal gate", ACCENT),
        ("L7", "shield", ACCENT), ("L8", "fail-safe", ACCENT),
        ("L9", "arbitration", ACCENT),
    ]):
        yy = y0 + i * Inches(0.34)
        rect(s, Inches(7.25), yy, Inches(2.7), Inches(0.30), fill=col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t2 = s.shapes[-1].text_frame
        t2.margin_left = Inches(0.1)
        p(t2, f"{code}   {name}", 10, WHITE, bold=True, first=True, space=0)
    # arrow to self-trust plane
    p(tb(s, Inches(10.02), Inches(3.7), Inches(0.5), Inches(0.4)),
      "→", 22, MUTE, first=True, space=0, align=PP_ALIGN.CENTER)
    rect(s, Inches(10.55), Inches(2.5), Inches(2.0), Inches(2.5), fill=INPT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tb(s, Inches(10.62), Inches(2.6), Inches(1.85), Inches(2.35))
    p(tf, "SELF-TRUST\nPLANE", 12, WHITE, bold=True, first=True, space=4,
      align=PP_ALIGN.CENTER, font=HEAD)
    p(tf, "measures each monitor's dynamic range, drift, and phase state", 9,
      RGBColor(0xE4, 0xD9, 0xEF), space=8, align=PP_ALIGN.CENTER)
    p(tf, "attaches a trust score to every verdict", 9, RGBColor(0xE4, 0xD9, 0xEF),
      space=0, align=PP_ALIGN.CENTER)
    foot(s, "The stack (L1-L9) is implemented and characterised. The self-trust plane is a design "
            "proposal justified by the current experimental findings.")

    # 8 SYSTEM ARCHITECTURE ------------------------------------------------
    s = slide(prs, "System architecture", "control and data flow")
    p(tb(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.3)),
      "Sensor input → estimation → trust → learned proposal → twin prediction → statistical "
      "and deterministic gates → fail-safe → arbitration → actuation.", 10.5, MUTE,
      first=True, space=0)

    # Top: sensors row
    rect(s, Inches(0.6), Inches(1.85), Inches(2.7), Inches(0.4), fill=MUTE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p(tb(s, Inches(0.6), Inches(1.92), Inches(2.7), Inches(0.28)),
      "SENSORS  IMU · GPS · LIDAR · CAN · fused frame", 9.5, WHITE, bold=True,
      first=True, space=0, align=PP_ALIGN.CENTER)

    # SHARED
    rect(s, Inches(0.6), Inches(2.45), Inches(4.6), Inches(2.2), fill=TINT2)
    p(tb(s, Inches(0.7), Inches(2.5), Inches(4.4), Inches(0.28)),
      "SHARED  (perception & estimation)", 10, INK, bold=True, first=True, space=0)
    for i, (c, txt) in enumerate([("L1", "Sensor bus, staleness rule, per-channel health"),
                                  ("L2a", "UKF innovation (Mahalanobis distance)"),
                                  ("L2b", "Fast-state estimate  ŷ, v̂, â"),
                                  ("L3", "Trust index (Mondrian conformal context)")]):
        yy = Inches(2.8) + i * Inches(0.42)
        rect(s, Inches(0.7), yy, Inches(0.6), Inches(0.34), fill=ACCENT_DK,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        p(tb(s, Inches(0.7), yy + Inches(0.05), Inches(0.6), Inches(0.24)),
          c, 10, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
        p(tb(s, Inches(1.35), yy + Inches(0.08), Inches(3.8), Inches(0.28)),
          txt, 9.5, BODY, first=True, space=0)

    # CORE-A
    rect(s, Inches(5.35), Inches(2.45), Inches(1.8), Inches(2.2), fill=RGBColor(0xEC, 0xEA, 0xF3))
    p(tb(s, Inches(5.45), Inches(2.5), Inches(1.6), Inches(0.28)),
      "CORE-A", 10, INK, bold=True, first=True, space=0)
    rect(s, Inches(5.45), Inches(2.9), Inches(0.6), Inches(0.34), fill=MUTE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p(tb(s, Inches(5.45), Inches(2.95), Inches(0.6), Inches(0.24)),
      "L4", 10, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
    p(tb(s, Inches(5.45), Inches(3.35), Inches(1.6), Inches(1.25)),
      "Learned proposer\n(PPO with PID-Lagrangian CMDP)\n\nWrites π_prop to Core-B; may not "
      "read Core-B artefacts.", 8.5, BODY, first=True, space=0, italic=True)

    # CORE-B
    rect(s, Inches(7.30), Inches(2.45), Inches(5.4), Inches(2.2), fill=TINT2)
    p(tb(s, Inches(7.4), Inches(2.5), Inches(5.2), Inches(0.28)),
      "CORE-B  (safety gates, verdict on every tick)", 10, INK, bold=True, first=True, space=0)
    for i, (c, txt) in enumerate([
        ("L5", "PINN twin — physics-based command prediction"),
        ("L6", "Statistical (conformal) gate — non-conformity score"),
        ("L7", "Deterministic shield + physical envelope"),
        ("L8", "Fail-safe FSM"),
    ]):
        yy = Inches(2.8) + i * Inches(0.42)
        rect(s, Inches(7.4), yy, Inches(0.6), Inches(0.34), fill=ACCENT,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        p(tb(s, Inches(7.4), yy + Inches(0.05), Inches(0.6), Inches(0.24)),
          c, 10, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
        p(tb(s, Inches(8.05), yy + Inches(0.08), Inches(4.55), Inches(0.28)),
          txt, 9.5, BODY, first=True, space=0)

    # Bottom: L9 + actuation
    rect(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(0.4), fill=ACCENT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p(tb(s, Inches(0.6), Inches(4.92), Inches(12.1), Inches(0.28)),
      "L9  ARBITRATION  →  actuation  (issues the safe command; records the audit-trail record)",
      10, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)

    # Observability
    rect(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.2), fill=RGBColor(0xF2, 0xEC, 0xF6))
    rect(s, Inches(0.6), Inches(5.4), Inches(0.14), Inches(1.2), fill=INPT)
    tf = tb(s, Inches(0.9), Inches(5.5), Inches(11.7), Inches(1.05))
    p(tf, "OBSERVABILITY PLANE  (offline; read-only taps into every layer above)", 10,
      INPT, bold=True, first=True, space=4)
    p(tf, "Fault-evidence profiler  ·  stability analyzer (σ_between, n_eff, autocorrelation)  ·  "
          "monitorability M(f,s)  ·  fault-evidence ledger  ·  redundancy-consistency monitor",
      10, BODY, space=3)
    p(tf, "The observability plane never writes to Core-A. It is a monitor for the monitors, "
          "and it is the substrate for the self-trust mechanism in ASTRA 2.0.",
      10, MUTE, italic=True, space=0)
    foot(s, "Diagram based on `src/astra/` module layout; 9/9 layers domain-mapped; 3/3 Core-B "
            "gates verdict on every one of 300 ticks in a live verification run.")

    # 9 WORKFLOW / METHODOLOGY --------------------------------------------
    s = slide(prs, "Workflow and methodology", "how a control tick is processed")
    tf = tb(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.3))
    p(tf, "Every 50 ms (20 Hz fast rate), the following sequence executes for one control tick.",
      10.5, MUTE, first=True, space=0)

    steps = [
        ("1", "Sensor acquisition", "IMU/GPS/LIDAR/CAN publish into a shared bus. Timestamps "
         "and per-channel health are recorded. Redundant sensing draws independent noise per "
         "channel from a disjoint seed."),
        ("2", "L1 fusion & health check", "Frame assembled; L1 evaluates staleness against a "
         "50 ms budget and flags per-modality health."),
        ("3", "L2 estimation (dual-rate UKF)", "Predict + update on the fused frame. Records "
         "the innovation vector, its Mahalanobis distance, and the covariance-normalised residual."),
        ("4", "L3 trust classification", "Context classifier (Mondrian scheme) assigns the tick "
         "to a context class (URBAN_CLEAR, HIGHWAY_CLEAR, DEGRADED_SENSOR); trust index emitted."),
        ("5", "L4 (Core-A) proposes π_prop", "Learned PPO policy consumes the fast state and "
         "trust and emits a command proposal. This is the only write from Core-A into Core-B."),
        ("6", "L5 twin prediction", "PINN twin produces the command that modelled physics "
         "expects, conditioned on the L3 context class."),
        ("7", "L6 statistical gate", "Compute non-conformity score = ‖π_prop − π_twin‖ / √P̂_a. "
         "Compare against the calibrated conformal quantile → PASS / VETO."),
        ("8", "L7 shield + L7b physical envelope", "Deterministic bounds on speed, lateral "
         "acceleration, jerk. VETO if any bound is exceeded."),
        ("9", "L8 fail-safe FSM", "Selects between issued command, degraded command, or halt "
         "posture based on gate verdicts and posture history."),
        ("10", "L9 arbitration → actuation", "Emits the safe command to the actuator; writes "
         "the full audit-trail record (fast state, proposal, gate verdicts, posture, seed)."),
    ]
    x0, y0 = Inches(0.6), Inches(1.85)
    for i, (n, ttl, body) in enumerate(steps):
        col_i = i % 2
        row_i = i // 2
        x = x0 + col_i * Inches(6.1)
        y = y0 + row_i * Inches(0.98)
        rect(s, x, y, Inches(6.0), Inches(0.90), fill=TINT if row_i % 2 else TINT2)
        rect(s, x, y, Inches(0.5), Inches(0.90), fill=ACCENT)
        p(tb(s, x, y + Inches(0.27), Inches(0.5), Inches(0.4)),
          n, 15, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER, font=HEAD)
        tf = tb(s, x + Inches(0.6), y + Inches(0.08), Inches(5.35), Inches(0.8))
        p(tf, ttl, 10.5, INK, bold=True, first=True, space=2, font=HEAD)
        p(tf, body, 9, BODY, space=0)
    foot(s, "Fault injection, when active, is applied at step 1/2 (sensor level) or step 5-6 "
            "(via the per-channel redundant sensing path). The rest of the pipeline is invariant.")

    # 10 TECHNOLOGY STACK -------------------------------------------------
    s = slide(prs, "Technology stack", "chosen for reproducibility and control")
    stacks = [
        ("Language & runtime", ACCENT, [
            ("Python 3.12", "single language across pipeline, training and analysis; "
             "reproducibility trumps performance at this stage"),
            ("PowerShell / Bash", "orchestration; the project is CPU-only Windows 11")]),
        ("Numerical & ML", ACCENT_DK, [
            ("NumPy 2.5.1", "primary numeric library; version pinned"),
            ("PyTorch 2.13.0", "neural policy and PINN twin; CPU inference"),
            ("Stable-Baselines3 2.9.0", "PPO implementation for the L4 policy"),
            ("Gymnasium", "training environment; NEVER imported by src/astra runtime")]),
        ("Statistics", ACCENT, [
            ("statistics.NormalDist + custom impl", "Wilcoxon, Holm-Bonferroni, BCa "
             "bootstrap, Spearman — hand-implemented to avoid mutating the runtime lockfile "
             "with scipy")]),
        ("Test & CI", ACCENT_DK, [
            ("pytest", "unit + integration; 82 passed / 2 failing guard tests (pre-existing)"),
            ("import-linter (contracts)", "enforces Core-A → Core-B one-way boundary as "
             "architectural invariant")]),
        ("Documentation & artefacts", ACCENT, [
            ("python-pptx (isolated venv)", "presentation generation without touching the "
             "measurement lockfile"),
            ("Markdown", "every experiment carries protocol, integrity checks, decision log")]),
        ("Version control", ACCENT_DK, [
            ("git", "branch `3.0`, six commits recording E17/E18/R1/R2/R3/R3b/R3c"),
        ]),
    ]
    x0, y0 = Inches(0.6), Inches(1.45)
    for i, (title, col, items) in enumerate(stacks):
        col_i = i % 2
        row_i = i // 2
        x = x0 + col_i * Inches(6.1)
        y = y0 + row_i * Inches(1.75)
        rect(s, x, y, Inches(6.0), Inches(1.68), fill=TINT if row_i % 2 else TINT2)
        rect(s, x, y, Inches(6.0), Inches(0.32), fill=col)
        p(tb(s, x + Inches(0.15), y + Inches(0.06), Inches(5.7), Inches(0.24)),
          title, 11, WHITE, bold=True, first=True, space=0, font=HEAD)
        tf = tb(s, x + Inches(0.2), y + Inches(0.4), Inches(5.6), Inches(1.25))
        for it, rea in items:
            p(tf, it, 10, INK, bold=True, space=2)
            p(tf, "why: " + rea, 9, BODY, space=5, italic=True)
    foot(s, "Lockfile discipline is deliberate: adding scipy/matplotlib to the measurement venv "
            "would put every prior latency and gate measurement in question.")

    # 11 CORE ALGORITHMS --------------------------------------------------
    s = slide(prs, "Core algorithms and technical approach", "the three that carry the paper")

    # left column - three algorithms
    algos = [
        ("L2  UKF innovation", ACCENT_DK,
         "What it does", "Unscented Kalman filter on the fast state; emits the innovation "
         "vector and its Mahalanobis distance under the innovation covariance.",
         "Why chosen", "Standard, well-understood non-linear estimator; innovation is the "
         "natural residual for FDI at the estimation stage.",
         "Inputs / outputs", "Fused sensor frame → posterior mean/covariance and residual "
         "statistics.",
         "Notes", "The innovation is the substrate for the L3 trust index and, in this study, "
         "for the D_L2a stage-wise discriminability."),
        ("L4  PPO with PID-Lagrangian CMDP", MUTE,
         "What it does", "Learns a lateral-control policy under three constraint costs "
         "(comfort, envelope, actuator).",
         "Why chosen", "PID-Lagrangian is the standard treatment for constrained MDPs; PPO is "
         "sample-efficient and well-supported.",
         "Training", "Offline, on the synthetic driving environment (kinematic bicycle). "
         "Multipliers held fixed within a batch, updated between batches.",
         "Notes", "This is Core-A. Its output is the untrusted proposal π_prop. It is NEVER "
         "trained on real vehicle data in the current work."),
        ("L6  Inductive Conformal Prediction", ACCENT,
         "What it does", "Non-parametric anomaly gate: computes non-conformity score "
         "s = ‖π_prop − π_twin‖ / √P̂_a and thresholds against a calibration corpus quantile.",
         "Why chosen", "Distribution-free guarantee of false-alarm rate under exchangeability; "
         "no parametric assumption on the residual distribution.",
         "Inputs / outputs", "Proposal and twin prediction → score → PASS / VETO verdict.",
         "Notes", "Whether exchangeability actually holds is the whole subject of Phase 5 (E18 "
         "series). It didn't with the legacy corpus; policy-conditional recalibration restored it."),
    ]
    y = Inches(1.4)
    for title, col, l1, t1, l2, t2, l3, t3, l4, t4 in algos:
        rect(s, Inches(0.6), y, Inches(8.4), Inches(1.75), fill=TINT)
        rect(s, Inches(0.6), y, Inches(0.14), Inches(1.75), fill=col)
        tf = tb(s, Inches(0.9), y + Inches(0.1), Inches(7.9), Inches(1.6))
        p(tf, title, 12.5, INK, bold=True, first=True, space=4, font=HEAD)
        for lbl, txt in [(l1, t1), (l2, t2), (l3, t3), (l4, t4)]:
            p(tf, f"{lbl}  ·  {txt}", 9.5, BODY, space=2)
        y += Inches(1.83)

    # right column - metric definitions
    rect(s, Inches(9.15), Inches(1.4), Inches(3.55), Inches(5.4), fill=TINT2)
    tf = tb(s, Inches(9.3), Inches(1.5), Inches(3.35), Inches(5.2))
    p(tf, "PRIMARY EVALUATION METRICS", 10, ACCENT, bold=True, first=True, space=8)
    p(tf, "D_s(f) — stage discriminability", 10.5, INK, bold=True, space=2)
    p(tf, "AUC( T_s | faulted , T_s | matched-clean ), folded to [0.5, 1.0]. "
          "0.5 = chance separability, 1.0 = perfect. Ticks pooled within a run.",
      9, BODY, space=8, italic=True)
    p(tf, "Run-level false-alarm rate", 10.5, INK, bold=True, space=2)
    p(tf, "|{t : score_t > q}| / |E_r| computed per run; unit of analysis is the run, "
          "never the tick.", 9, BODY, space=8, italic=True)
    p(tf, "Per-run FAR band [ε/2, 2ε]", 10.5, INK, bold=True, space=2)
    p(tf, "At ε=0.05, band = [2.5%, 10%]. A calibration is 'valid' only if the majority of "
          "held-out clean runs fall inside this band.", 9, BODY, space=8, italic=True)
    p(tf, "Phase-resolved detection", 10.5, INK, bold=True, space=2)
    p(tf, "During-fault and post-fault alarm rates reported separately. Aggregate rates would "
          "have hidden the R3b/R3c finding.", 9, BODY, space=0, italic=True)

    # 12 IMPLEMENTATION --------------------------------------------------
    s = slide(prs, "Implementation", "modules and integration")
    tf = tb(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.5))
    heading(tf, "Repository structure (implemented)", first=True)
    para(tf, [
        "src/astra/  — nine-layer runtime governance stack (production code, no training deps)",
        "training/  — offline training for L4 policy and L5 twin; calibration harvest for L6",
        "benchmarks/  — experiment drivers (E17, E18, R1, R2, R3, R3b, R3c) and analysis",
        "tests/  — unit + integration + guard tests, contract enforcement",
        "research/  — experiment history, audits, invalidation records, claim ledger",
        "experiments/phase5_od8_h7/  — pre-registrations, frozen thresholds, decision logs",
        "results/  and  var/policy/, var/twin/, var/calibration/  — artefacts, checkpoints",
    ], size=10, space=3)

    heading(tf, "Key implementation decisions")
    para(tf, [
        "Runtime and training are import-separated: nothing in src/astra imports Gymnasium/SB3.",
        "Every experiment is pre-registered before running; the criterion is frozen on disk.",
        "Delivered-signal integrity check is a workflow requirement, not a coding convention — "
        "results without it are marked invalid.",
        "Invalidated results are retained with a WITHDRAWN flag; nothing is silently deleted.",
        "Statistical library is hand-implemented in numpy + stdlib to avoid mutating the "
        "measurement lockfile.",
    ], size=10, space=3)

    rect(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.5), fill=TINT)
    tf = tb(s, Inches(7.0), Inches(1.55), Inches(5.6), Inches(5.2))
    p(tf, "MODULE INTERACTION (SIMPLIFIED)", 10, ACCENT, bold=True, first=True, space=6)
    p(tf, "training.closed_loop.drive_closed_loop( )", 10, INK, bold=True, space=2, font="Consolas")
    p(tf, "→ constructs the pipeline, drives it for N ticks, calls an observer per tick, "
          "returns aggregate results", 9, BODY, space=8, italic=True)
    p(tf, "benchmarks.e18r3b_detect  (example)", 10, INK, bold=True, space=2, font="Consolas")
    p(tf, "1. Load frozen threshold (v3, hard-coded)\n"
          "2. For each fault × seed, run one clean + one faulted arm\n"
          "3. Record per-tick scores; integrity-check delivered signal\n"
          "4. Compute per-run alarm rate vs run-level decision boundary",
      9, BODY, space=10)
    p(tf, "SCREENSHOTS OF THE APPLICATION UI", 10, INPT, bold=True, space=4)
    p(tf, "[NEEDS INPUT] — the project has no user-facing UI; command-line drivers and "
          "generated PowerPoint / Markdown reports are the primary interfaces. Suggest: "
          "insert screenshots of the terminal running an experiment, and of a generated "
          "figure such as fig1_stage_profile_P1.svg.", 9, MUTE, italic=True, space=0)
    foot(s, "Runtime code path is 21,631 SLOC; test code is 29,343 SLOC (>1:1 test-to-code ratio).")

    # 13 CURRENT PROJECT STATUS -------------------------------------------
    s = slide(prs, "Current project status", "what is complete and what is in flight")

    # left: DONE
    rect(s, Inches(0.6), Inches(1.4), Inches(6.05), Inches(5.5), fill=RGBColor(0xE7, 0xF1, 0xEA))
    rect(s, Inches(0.6), Inches(1.4), Inches(0.14), Inches(5.5), fill=OK)
    tf = tb(s, Inches(0.9), Inches(1.5), Inches(5.6), Inches(5.3))
    p(tf, "COMPLETED", 11, OK, bold=True, first=True, space=6)
    p(tf, "Architecture and runtime", 10.5, INK, bold=True, space=2, font=HEAD)
    para(tf, [
        "All 9 layers implemented and domain-mapped",
        "All 3 Core-B gates verdict on every tick (verified live on 300 ticks)",
        "Core-A → Core-B one-way boundary enforced by contract tests",
    ], size=9.5, space=2)
    p(tf, "Training pipeline", 10.5, INK, bold=True, space=2, font=HEAD)
    para(tf, [
        "L4 policy trained (three checkpoints: synthetic, long, jerkscaled)",
        "L5 PINN twin trained on kinematic bicycle dynamics",
        "L6 calibration corpus harvested",
    ], size=9.5, space=2)
    p(tf, "Experimental campaigns", 10.5, INK, bold=True, space=2, font=HEAD)
    para(tf, [
        "E17 30-seed sweep: 90 profiles, 2,160 runs",
        "E17 position re-injection: 720 runs (invalidated original claim)",
        "E18 calibration series: 72,000 clean ticks + 1,260 faulted runs",
        "E18-R1 (run-local), R2 (matched-pooled), R3 (long window), R3b (detection), "
        "R3c (sustained fault): together +720 runs",
        "Zero execution failures across ~4,300 runs",
    ], size=9.5, space=2)

    # right: IN PROGRESS
    rect(s, Inches(6.85), Inches(1.4), Inches(5.85), Inches(5.5), fill=RGBColor(0xF7, 0xEE, 0xDA))
    rect(s, Inches(6.85), Inches(1.4), Inches(0.14), Inches(5.5), fill=WARN)
    tf = tb(s, Inches(7.15), Inches(1.5), Inches(5.4), Inches(5.3))
    p(tf, "CURRENTLY IN PROGRESS / NEXT", 11, WARN, bold=True, first=True, space=6)
    p(tf, "Systematic literature review", 10.5, INK, bold=True, space=2, font=HEAD)
    p(tf, "Cross-check that the self-trust / monitorability contribution is unclaimed. "
          "In progress; blocks the novelty claim.", 9.5, BODY, space=6)
    p(tf, "Phase 2 — Monitorability re-analysis", 10.5, INK, bold=True, space=2, font=HEAD)
    p(tf, "Compute M(f,s) on the 28 existing cells. Zero new compute, hours of work; validates "
          "or removes the paper's proposed metric.", 9.5, BODY, space=6)
    p(tf, "P3 recalibration (E18-R4)", 10.5, INK, bold=True, space=2, font=HEAD)
    p(tf, "Fix P3 threshold bias identified in R3; short compute.", 9.5, BODY, space=6)
    p(tf, "E19 — H7 monitor placement", 10.5, INK, bold=True, space=2, font=HEAD)
    p(tf, "Unblocked on P1 after E18-R3 PASS; awaits Phase 2 result before design finalisation.",
      9.5, BODY, space=6)
    p(tf, "Self-trust layer implementation", 10.5, INK, bold=True, space=2, font=HEAD)
    p(tf, "Design complete (ASTRA 2.0 review). Implementation as production code is the "
          "next major engineering task.", 9.5, BODY, space=0)

    # 14 ROADMAP -----------------------------------------------------------
    s = slide(prs, "Roadmap  —  current stage to final demo", "where the project is going")
    stages = [
        ("Completed", "5 experiment phases; 4,300 runs; audited findings; withdrawn claims "
         "documented; ASTRA 2.0 architecture proposed", DONE_TXT := "DONE", OK),
        ("Current stage", "Literature review + Phase 2 monitorability re-analysis + P3 "
         "recalibration + E19 pre-registration", "IN PROGRESS", WARN),
        ("Next steps", "Implement self-trust layer + phase-aware detector; complete E19 monitor "
         "placement study; run E20 lying-sensor extension on the per-channel injection path", "NEXT",
         ACCENT),
        ("Final implementation", "Integrated ASTRA 2.0 with observability plane, monitorability "
         "output, and phase-resolved decision rule", "PLANNED", MUTE),
        ("Testing & validation", "comma2k19 real-driving-data replay (open-loop); CARLA closed-"
         "loop only if time permits; baseline comparison against innovation-χ² monitor", "PLANNED", MUTE),
        ("Final demo", "Live end-to-end: inject faults into a driven scenario, show phase-"
         "resolved detection, show the self-trust plane annotating each verdict", "PLANNED", MUTE),
    ]
    y = Inches(1.5)
    for i, (title, body, tag, col) in enumerate(stages):
        rect(s, Inches(0.6), y, Inches(12.1), Inches(0.83), fill=TINT if i % 2 else TINT2)
        rect(s, Inches(0.6), y, Inches(0.16), Inches(0.83), fill=col)
        # stage number
        p(tb(s, Inches(0.85), y + Inches(0.24), Inches(0.7), Inches(0.4)),
          f"0{i + 1}", 22, col, bold=True, first=True, space=0, font=HEAD)
        # title + body
        tf = tb(s, Inches(1.7), y + Inches(0.08), Inches(9.4), Inches(0.72))
        p(tf, title, 12, INK, bold=True, first=True, space=3, font=HEAD)
        p(tf, body, 10, BODY, space=0)
        # chip
        chip(s, Inches(11.35), y + Inches(0.29), tag)
        y += Inches(0.87)
    foot(s, "Dates are deliberately not specified. Progression is gated: a stage is not entered "
            "just because it exists in the plan. External datasets follow, not lead, the internal "
            "causal chain.")

    # 15 RESULTS ----------------------------------------------------------
    s = slide(prs, "Results and preliminary findings", "measured, not assumed")

    tf = tb(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.32))
    p(tf, "All results below are on the synthetic plant with three learned policy checkpoints. "
          "Positive findings are supported by ≥ 30 seeds and pre-registered criteria. "
          "External-data validation is future work.", 10, MUTE, first=True, space=0)

    # positive findings
    table(s, Inches(0.6), Inches(1.9), Inches(12.1), [
        ["Finding", "Evidence", "Status"],
        ["OD-8 miscalibration is a calibration-set provenance failure, not a threshold-value one",
         "Global recalibration reproduces the original defect exactly (0.00 % FAR on P1/P3, "
         "11.06 % on P2)", "Established"],
        ["P1 conformal monitor: stable per-run false-alarm rate at a 160-second window",
         "30/30 held-out clean runs in the [2.5 %, 10 %] band at n = 3,200 (E18-R3)", "Established"],
        ["Windowed calibration limit was precision, not dynamics",
         "Per-run FAR variance scales with n; slope −0.739 (P1), −0.918 (P3) (E18-R3)", "Established"],
        ["4 of 6 faults reach 100 % detection at n = 3,200",
         "position_bias, position_drift, lateral_noise, imu_dropout at 100 % (E18-R3b)",
         "Established (with qualifier below)"],
        ["A persistent sensor failure is invisible to the monitor for as long as it persists",
         "Sustained imu_dropout, 160 s, 30 seeds: alarm rate 0.2 % — BELOW clean baseline of "
         "~5 % (E18-R3c). R3b's 100 % detection was the post-fault recovery transient.",
         "Established"],
        ["Statistical discriminability does not predict operational detection",
         "17 of 28 cells disagree between D_L6 ≥ 0.9 and detection ≥ 0.9. D_L1 vs detection: "
         "Spearman ρ = −0.480, p = 0.0088 (E18)", "Established"],
        ["Fault-induced alarm suppression is a real phenomenon, not an anecdote",
         "11 of 28 cells, all p < 0.05. imu_dropout on P1: 0.10 % alarm during fault vs 5.47 % "
         "clean baseline (E18); phase-resolved mechanism confirmed by R3c", "Established"],
        ["Redundancy can conceal a fault entirely",
         "POSITION_Y was inert against the driven sensing path for 17 days (E17-Position); "
         "Control C established causally: estimator sees 0 with redundancy on, 3.18 with it off",
         "Established"],
    ], widths=[3.5, 6.7, 1.9], fs=9, hfs=9.5, rh=Inches(0.6))
    foot(s, "Six earlier claims were withdrawn during the campaign after measurement defects were "
            "self-caught. All are documented in the project claim ledger; the discipline of "
            "recording them is itself part of the contribution.")

    # 16 CHALLENGES & LIMITATIONS ----------------------------------------
    s = slide(prs, "Challenges and limitations", "honestly stated")

    tf = tb(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.5))
    heading(tf, "Technical challenges encountered", first=True)
    para(tf, [
        "Six measurement defects self-caught during the campaign, five before publication. "
        "Each withdrew a claim and each was recorded rather than hidden.",
        "Original position-fault result was 0.5-by-construction: the redundant sensing path "
        "was regenerating the position channel from ground truth. Corrected by re-injecting "
        "through a per-channel path.",
        "Original P1 'VALID' calibration was measured on a different tick window than "
        "detection. Withdrawn once matched-window analysis exposed it.",
        "R3b's 100 % detection of imu_dropout was traced to the post-fault recovery transient, "
        "not the fault itself. R3c (duration-matched control) settled this.",
    ], size=10, space=3)

    heading(tf, "How the team addressed them")
    para(tf, [
        "Adopted pre-registration: criterion frozen on disk before each experiment ran.",
        "Adopted delivered-signal integrity as a workflow requirement, not a coding convention.",
        "Adopted the run as the statistical unit; tick-level rates are now supplementary only.",
        "Maintain a machine-readable claim ledger; nothing is silently deleted.",
    ], size=10, space=3)

    tf = tb(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.5))
    heading(tf, "Known limitations of the current work", first=True)
    para(tf, [
        "Synthetic plant only. All findings are [M-syn]; external validation [M-ext] is 0 of 30. "
        "This is the largest single limitation and the reason external datasets are the next stage.",
        "Three policies is not a sample of policies; findings on P1 do not automatically "
        "generalise across policy space.",
        "P2 is INVALID for fixed-quantile calibration (non-stationary score). Retained in the "
        "record; not tuned around.",
        "Two faults (speed_bias, speed_stuck) remain undetectable by the current L6 gate at any "
        "tested severity; their non-conformity score is flat in every phase.",
        "No end-to-end user-facing application; the deliverable is an audited testbed, "
        "measurement discipline, findings, and an architecture proposal.",
        "Self-trust plane is a design; implementation and validation are future work.",
        "Novelty of the specific self-trust framing awaits systematic literature review.",
    ], size=10, space=3)
    foot(s, "The withdrawal record is treated as an asset of the project: a reviewer can inspect "
            "exactly what was tried, what failed, and why, rather than only the polished result.")

    # 17 FUTURE SCOPE ----------------------------------------------------
    s = slide(prs, "Future scope", "realistic extensions")
    items = [
        ("Near-term (weeks)",
         "Complete literature review (blocks the novelty claim); compute monitorability "
         "metric on the 28 existing cells (zero new compute); recalibrate P3 (short compute); "
         "run E19 monitor-placement experiment on P1 with pre-registered baselines.", ACCENT),
        ("Medium-term (months)",
         "Implement the self-trust plane as production code with an interface, not a "
         "benchmark script. Run E20 (lying-sensor / adversarial): one channel lies within "
         "the voting margin using the redundancy-consistency monitor as the detection "
         "substrate.", ACCENT_DK),
        ("External validation",
         "comma2k19: open-loop replay of real driving data with controlled fault injection "
         "layered on top. Requires an adapter and confirmed timestamp precision.",
         MUTE),
        ("Naturalistic regime validation",
         "highD: naturalistic traffic trajectories to test whether findings survive across "
         "operating regimes rather than one synthetic driving profile.",
         MUTE),
        ("Closed-loop realism",
         "CARLA: closed-loop simulation with better sensor models; useful demonstration but "
         "not a substitute for real data.",
         MUTE),
        ("Standards and industry engagement",
         "Position monitorability + self-trust as a SOTIF (ISO 21448) coverage contribution; "
         "engage automotive R&D partners for real sensor data via NDA. Explicitly out of "
         "scope: any autonomy claim beyond fault detection and monitor assurance.",
         BAD),
    ]
    y = Inches(1.45)
    for title, body, col in items:
        rect(s, Inches(0.6), y, Inches(12.1), Inches(0.87), fill=TINT)
        rect(s, Inches(0.6), y, Inches(0.16), Inches(0.87), fill=col)
        tf = tb(s, Inches(0.95), y + Inches(0.1), Inches(11.5), Inches(0.72))
        p(tf, title, 11, INK, bold=True, first=True, space=3, font=HEAD)
        p(tf, body, 10, BODY, space=0)
        y += Inches(0.92)

    # 18 CONCLUSION ------------------------------------------------------
    s = slide(prs, "Conclusion", "in summary", dark=True)
    rect(s, 0, 0, W, H, fill=HEAD_BG)
    p(tb(s, Inches(0.6), Inches(0.38), Inches(12.1), Inches(0.24)),
      "IN SUMMARY", 9.5, ACCENT, bold=True, first=True, space=0)
    p(tb(s, Inches(0.6), Inches(0.64), Inches(12.1), Inches(0.6)),
      "Conclusion", 26, WHITE, bold=True, first=True, space=0, font=HEAD)

    blocks = [
        ("Problem",
         "Autonomous systems with learned controllers rely on runtime monitors that existing "
         "architectures assume are reliable. That assumption is not measured."),
        ("Proposed solution",
         "ASTRA: a nine-layer runtime governance stack (Simplex-derived) instrumented so fault "
         "evidence is measurable at every stage, plus an observability plane that assesses "
         "whether each monitor can be trusted."),
        ("Current achievement",
         "Six audited experiment phases, ~4,300 runs, zero fabricated results, six withdrawn "
         "claims documented. A working per-run-stable monitor on one policy at a 160-second "
         "window, four of six faults detectable, and — the most significant single finding — "
         "measured evidence that the monitor is silent for the duration of a persistent "
         "sensor failure."),
        ("Expected final outcome",
         "An implemented self-trust plane, monitorability validated as a predictor of "
         "operational detection, and an initial external-data replication on comma2k19. "
         "Publishable at an IEEE conference in the runtime-assurance / dependability track."),
        ("Overall significance",
         "The contribution is a shift from 'is there a fault?' to 'can the monitor answering "
         "that question currently be trusted?' — a distinction our own experiments show is "
         "necessary, and one no existing runtime-assurance architecture makes."),
    ]
    y = Inches(1.4)
    for ttl, body in blocks:
        rect(s, Inches(0.6), y, Inches(12.1), Inches(1.02), fill=RGBColor(0x1E, 0x36, 0x42))
        rect(s, Inches(0.6), y, Inches(0.16), Inches(1.02), fill=ACCENT)
        tf = tb(s, Inches(0.95), y + Inches(0.1), Inches(11.5), Inches(0.88))
        p(tf, ttl, 12, ACCENT, bold=True, first=True, space=4, font=HEAD)
        p(tf, body, 11, RGBColor(0xD2, 0xE1, 0xE7), space=0)
        y += Inches(1.08)

    prs.save(str(out))
    print(f"  {len(prs.slides._sldIdLst)} slides -> {out}")


if __name__ == "__main__":
    build(Path(__file__).resolve().parent / "ASTRA_Panel_Review.pptx")

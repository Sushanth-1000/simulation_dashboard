"""ASTRA 12-slide research status deck.

Reflects the state after E18-R2 (1 September 2026). Every number is traceable to
`research/`, `experiments/phase5_od8_h7/`, or `results/`. Claim wording follows
`CLAIM_LEDGER.md` — in particular, no OD-8 operational-monitoring claim is made.

Run with the isolated docs venv (python-pptx), NOT the measurement venv.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

W, H = Inches(13.333), Inches(7.5)

# Ocean Gradient — chosen for a systems/measurement subject, not generic blue.
DEEP = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
MID = RGBColor(0x21, 0x29, 0x5C)
INK = RGBColor(0x1A, 0x22, 0x33)
BODY = RGBColor(0x45, 0x55, 0x66)
FAINT = RGBColor(0x8A, 0x99, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xEF, 0xF4, 0xF8)
TINT2 = RGBColor(0xE4, 0xEC, 0xF2)

OK = RGBColor(0x1E, 0x6B, 0x52)
WARN = RGBColor(0xA9, 0x76, 0x0B)
BAD = RGBColor(0xA3, 0x33, 0x28)
GREY = RGBColor(0x6B, 0x72, 0x80)

CHIP = {"ESTABLISHED": OK, "SUPPORTED": OK, "COMPLETE": OK, "PARTIAL": WARN,
        "IN PROGRESS": WARN, "PRELIMINARY": WARN, "NEXT": TEAL, "PLANNED": GREY,
        "NOT STARTED": GREY, "BLOCKED": BAD, "WITHDRAWN": BAD, "REJECTED": BAD,
        "NOT TESTED": GREY, "HYPOTHESIS": GREY}

HEAD = "Cambria"
TEXT = "Calibri"


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    t = s.shapes.add_textbox(x, y, w, h).text_frame
    t.word_wrap = True
    t.vertical_anchor = anchor
    t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
    return t


def p(tf, text, size, color=INK, bold=False, first=False, space=6, font=TEXT,
      align=PP_ALIGN.LEFT, italic=False):
    par = tf.paragraphs[0] if first else tf.add_paragraph()
    par.alignment = align
    par.space_after = Pt(space)
    r = par.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return par


def rect(s, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE, line=None):
    o = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        o.fill.background()
    else:
        o.fill.solid()
        o.fill.fore_color.rgb = fill
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = line
        o.line.width = Pt(1)
    o.shadow.inherit = False
    return o


def chip(s, x, y, label, w=None):
    """The deck's repeating motif: a claim-status chip."""
    col = CHIP.get(label, GREY)
    w = w or Inches(0.16 + 0.088 * len(label))
    o = rect(s, x, y, w, Inches(0.25), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    t = o.text_frame
    t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
    t.word_wrap = False
    p(t, label, 8.5, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
    return w


def dot(s, x, y, n, col=DEEP, d=Inches(0.34)):
    """Numbered disc — used instead of bullets for step lists."""
    o = rect(s, x, y, d, d, fill=col, shape=MSO_SHAPE.OVAL)
    t = o.text_frame
    t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
    p(t, str(n), 11.5, WHITE, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)


def slide(prs, title=None, kicker=None, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = MID if dark else WHITE
    if title:
        if kicker:
            p(tb(s, Inches(0.7), Inches(0.40), Inches(11.9), Inches(0.26)),
              kicker.upper(), 10, TEAL, bold=True, first=True, space=0)
        p(tb(s, Inches(0.7), Inches(0.66), Inches(11.9), Inches(0.62)),
          title, 32, INK if not dark else WHITE, bold=True, first=True, space=0, font=HEAD)
    return s


def table(s, x, y, w, rows, widths=None, fs=10.5, hfs=10, rh=Inches(0.32),
          hh=Inches(0.34), colors=None, head=DEEP):
    t = s.shapes.add_table(len(rows), len(rows[0]), x, y, w, hh + rh * (len(rows) - 1)).table
    t.rows[0].height = hh
    for i in range(1, len(rows)):
        t.rows[i].height = rh
    if widths:
        tot = sum(widths)
        for j, fr in enumerate(widths):
            t.columns[j].width = Emu(int(w * fr / tot))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = head if i == 0 else (WHITE if i % 2 else TINT)
            tf = c.text_frame
            tf.word_wrap = True
            par = tf.paragraphs[0]
            par.space_after = Pt(0)
            r = par.add_run()
            r.text = str(val)
            r.font.size = Pt(hfs if i == 0 else fs)
            r.font.bold = i == 0
            r.font.name = TEXT
            if i == 0:
                r.font.color.rgb = WHITE
            elif colors and (i, j) in colors:
                r.font.color.rgb = colors[(i, j)]
                r.font.bold = True
            else:
                r.font.color.rgb = INK
    return t


def card(s, x, y, w, h, heading, lines, tint=TINT, hcol=DEEP, hs=12.5):
    rect(s, x, y, w, h, fill=tint)
    tf = tb(s, x + Inches(0.22), y + Inches(0.16), w - Inches(0.44), h - Inches(0.3))
    p(tf, heading, hs, hcol, bold=True, first=True, space=7, font=HEAD)
    for ln in lines:
        if isinstance(ln, tuple):
            p(tf, ln[0], 10.5, ln[1], bold=(len(ln) > 2 and ln[2]), space=5)
        else:
            p(tf, ln, 10.5, BODY, space=5)


def stat(s, x, y, w, value, label, col=DEEP, vs=34):
    # Box sized to actually contain headline + wrapped label, not just to look right.
    tf = tb(s, x, y, w, Inches(1.08))
    p(tf, value, vs, col, bold=True, first=True, space=1, font=HEAD, align=PP_ALIGN.CENTER)
    p(tf, label, 9.5, BODY, space=0, align=PP_ALIGN.CENTER)


def foot(s, text, col=FAINT):
    p(tb(s, Inches(0.7), Inches(6.96), Inches(11.9), Inches(0.3)),
      text, 8.5, col, first=True, space=0, italic=True)


def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---- 1 TITLE ---------------------------------------------------------
    s = slide(prs, dark=True)
    rect(s, 0, 0, W, H, fill=MID)
    rect(s, Inches(8.6), Inches(0), Inches(4.733), H, fill=RGBColor(0x1A, 0x21, 0x4B))
    for i, (v, l) in enumerate([("6", "experiment\nphases run"), ("4,290", "closed-loop\nruns executed"),
                                ("6", "measurement defects\nfound, 5 self-caught")]):
        yy = Inches(1.55) + i * Inches(1.62)
        tf = tb(s, Inches(9.15), yy, Inches(3.7), Inches(1.35))
        p(tf, v, 33, RGBColor(0x6E, 0xC8, 0xE8), bold=True, first=True, space=2, font=HEAD)
        for k, line in enumerate(l.split("\n")):
            p(tf, line, 10.5, RGBColor(0xA8, 0xC0, 0xD4), space=0)
    p(tb(s, Inches(0.85), Inches(1.5), Inches(7.4), Inches(0.3)),
      "RESEARCH STATUS BRIEFING", 11.5, RGBColor(0x6E, 0xC8, 0xE8), bold=True, first=True, space=0)
    tf = tb(s, Inches(0.85), Inches(2.05), Inches(7.4), Inches(3.0))
    p(tf, "ASTRA", 54, WHITE, bold=True, first=True, space=4, font=HEAD)
    p(tf, "Fault Observability and Operational", 21, RGBColor(0xD6, 0xE3, 0xEE), space=2, font=HEAD)
    p(tf, "Monitoring in a Layered Governance Stack", 21, RGBColor(0xD6, 0xE3, 0xEE), space=14, font=HEAD)
    p(tf, "Where does sensor-fault evidence stop being usable — and can that "
          "be turned into a working monitor?", 13, RGBColor(0x9F, 0xB8, 0xCC), space=0)
    p(tb(s, Inches(0.85), Inches(6.35), Inches(7.4), Inches(0.7)),
      "Sushanth C.  ·  Tarun Gowda V  ·  T Tilak Reddy\n"
      "Guide: Dr. Chaitra R.  ·  BMS College of Engineering  ·  1 September 2026",
      10.5, RGBColor(0x82, 0x9C, 0xB4), first=True, space=0)

    # ---- 2 INTRODUCTION --------------------------------------------------
    s = slide(prs, "The problem we are working on", "introduction")
    tf = tb(s, Inches(0.7), Inches(1.52), Inches(6.0), Inches(4.6))
    p(tf, "A sensor fault does not stay a sensor fault.", 15, INK, bold=True, first=True,
      space=9, font=HEAD)
    p(tf, "It passes through estimation, fusion, trust scoring and a stack of safety "
          "monitors. Each stage transforms it. Somewhere along that path the evidence "
          "that anything is wrong can become too weak for a downstream monitor to act on.",
      12, BODY, space=10)
    p(tf, "Conventional fault detection answers a binary question — “fault or no fault?” "
          "It does not tell an engineer WHERE in their own architecture the evidence "
          "stopped being usable, which is exactly what decides where a monitor should sit.",
      12, BODY, space=10)
    p(tf, "ASTRA is a 9-layer runtime governance stack for a learned controller, "
          "instrumented so that fault evidence can be measured at every stage.",
      12, INK, bold=True, space=0)

    card(s, Inches(7.05), Inches(1.52), Inches(5.55), Inches(2.05),
         "The primary question",
         [("When a sensor becomes faulty, where in the pipeline does the fault stop "
           "being distinguishable from healthy behaviour?", INK, True)], TINT, DEEP)
    card(s, Inches(7.05), Inches(3.77), Inches(5.55), Inches(1.28),
         "Secondary question",
         [("Can that information decide where monitoring should occur?", INK, True)],
         TINT2, TEAL)
    card(s, Inches(7.05), Inches(5.25), Inches(5.55), Inches(1.28),
         "Future security question",
         [("Does a deliberately manipulated sensor behave differently from a naturally "
           "faulty one?", GREY, True)], TINT2, GREY)
    chip(s, Inches(11.55), Inches(5.32), "NOT TESTED")
    foot(s, "Nine layers: L1 sensing · L2 dual-rate UKF · L3 trust · L4 learned proposer (Core-A) · "
            "L5 PINN twin · L6 conformal gate · L7 shield · L8 fail-safe · L9 arbitration.")

    # ---- 3 METHODOLOGY ---------------------------------------------------
    s = slide(prs, "How we measure", "methodology")
    p(tb(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.3)),
      "Four commitments that make the numbers mean something", 12.5, BODY, first=True, space=0)
    items = [
        ("Stage-wise discriminability", "D_s(f) = AUC(faulted, matched-clean), folded to [0.5, 1.0]. "
         "0.5 is chance-level separation, never “50 % accuracy”. Reported per stage as a profile, "
         "not collapsed to one number."),
        ("Delivered-signal integrity", "A result is rejected unless the faulted arm is shown to differ "
         "from the clean arm at the point each consumer reads it. This rule exists because three "
         "earlier results were 0.5 by construction."),
        ("Pre-registration", "Stage, faults, thresholds, tests and falsification criteria are written "
         "and frozen before execution. Deviations are reported as deviations."),
        ("The run is the unit", "Ticks inside a run are autocorrelated and are never treated as "
         "independent samples. Changing tick→run overturned two of our own verdicts."),
    ]
    for i, (h, body) in enumerate(items):
        x = Inches(0.7) + (i % 2) * Inches(6.05)
        y = Inches(2.0) + (i // 2) * Inches(2.15)
        rect(s, x, y, Inches(5.85), Inches(1.92), fill=TINT if i % 2 == 0 else TINT2)
        dot(s, x + Inches(0.24), y + Inches(0.22), i + 1)
        tf = tb(s, x + Inches(0.72), y + Inches(0.2), Inches(4.9), Inches(1.6))
        p(tf, h, 13, DEEP, bold=True, first=True, space=6, font=HEAD)
        p(tf, body, 10.5, BODY, space=0)
    foot(s, "Metric appropriateness is audited per fault type: dispersion faults use a rolling-standard-deviation "
            "statistic, dropout uses stream health, bias and drift use the raw channel.")

    # ---- 4 ARCHITECTURE --------------------------------------------------
    s = slide(prs, "The measurement surface", "architecture")
    p(tb(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.3)),
      "Verified in a live run, not read from documentation: 9 / 9 layers domain-mapped, "
      "3 / 3 Core-B gates return a verdict on every one of 300 ticks.", 11.5, BODY,
      first=True, space=0)
    cells = [("L1", "sensing", DEEP), ("L2a", "innovation", DEEP), ("L2b", "estimate", DEEP),
             ("L3", "trust", DEEP), ("L4", "proposer", GREY), ("L5", "twin", TEAL),
             ("L6", "conformal", WARN), ("L7", "shield", TEAL), ("L8", "fail-safe", TEAL),
             ("L9", "arbitration", GREY)]
    x = Inches(0.7)
    bw, gap = Inches(1.13), Inches(0.088)
    for code, name, col in cells:
        o = rect(s, x, Inches(2.08), bw, Inches(0.92), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t = o.text_frame
        t.word_wrap = True
        t.margin_left = t.margin_right = Inches(0.04)
        p(t, code, 14, WHITE, bold=True, first=True, space=1, align=PP_ALIGN.CENTER, font=HEAD)
        p(t, name, 8.5, RGBColor(0xDD, 0xE8, 0xF0), space=0, align=PP_ALIGN.CENTER)
        x += bw + gap
    for lo, wd, lab in [(Inches(0.7), Inches(4.72), "SHARED   L1–L3"),
                        (Inches(5.5), Inches(1.13), "CORE-A"),
                        (Inches(6.71), Inches(4.68), "CORE-B   3 gates, every tick")]:
        rect(s, lo, Inches(3.1), wd, Inches(0.3), fill=TINT2)
        p(tb(s, lo, Inches(3.14), wd, Inches(0.25)), lab, 9.5, BODY, first=True, space=0,
          align=PP_ALIGN.CENTER)
    table(s, Inches(0.7), Inches(3.68), Inches(11.9), [
        ["Stage", "Statistic measured", "What it tells us"],
        ["L1", "raw channel value / dispersion / stream health", "sensor-level detectability"],
        ["L2a", "UKF innovation (Mahalanobis distance)", "estimator-level observability"],
        ["L2b / L3", "estimated lateral position; trust index", "post-estimation state and trust"],
        ["L6", "conformal non-conformity score", "statistical gate evidence — and its threshold"],
        ["L7 / L8", "veto flag; fail-safe posture", "operational gate outcome"],
    ], widths=[1.2, 4.6, 4.4], fs=10.5)
    foot(s, "The one-way Core-A → Core-B trust boundary is enforced by architectural invariants: "
            "Core-A may write a proposal into Core-B and may not read any Core-B artefact.")

    # ---- 5 NOVELTY -------------------------------------------------------
    s = slide(prs, "What is actually new", "novelty")
    p(tb(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.3)),
      "No single ingredient is novel. The proposed contribution is the combination — and its "
      "status is tracked honestly, item by item.", 11.5, BODY, first=True, space=0)
    rows = [
        ("1", "Stage-wise fault discriminability across a governance pipeline", "ESTABLISHED"),
        ("2", "Fault-specific observability profiles — four distinct behaviours measured", "ESTABLISHED"),
        ("3", "A candidate absorption point A(f), reported only where well-posed", "PARTIAL"),
        ("4", "Separating statistical discriminability from operational gate detection", "ESTABLISHED"),
        ("5", "Diagnosing why a conformal monitor fails, at the level of its score process", "ESTABLISHED"),
        ("6", "Using measured observability to choose monitor placement (H7 / E19)", "BLOCKED"),
        ("7", "Natural vs adversarial sensor manipulation (E20)", "NOT TESTED"),
    ]
    y = Inches(2.0)
    for n, txt, st in rows:
        rect(s, Inches(0.7), y, Inches(11.9), Inches(0.56), fill=TINT if int(n) % 2 else TINT2)
        dot(s, Inches(0.86), y + Inches(0.11), n, col=DEEP if st in ("ESTABLISHED", "PARTIAL") else GREY)
        p(tb(s, Inches(1.42), y + Inches(0.17), Inches(8.9), Inches(0.3)), txt, 11.5, INK,
          first=True, space=0)
        chip(s, Inches(10.6), y + Inches(0.155), st)
        y += Inches(0.62)
    card(s, Inches(0.7), Inches(6.42), Inches(11.9), Inches(0.62),
         "", [("Novelty remains a research hypothesis until confirmed against a final systematic "
               "literature review and completed experiments. No priority claim is made.", INK, True)],
         TINT2, GREY, hs=1)

    # ---- 6 COMPARATIVE ---------------------------------------------------
    s = slide(prs, "How this differs from existing practice", "comparative analysis")
    table(s, Inches(0.7), Inches(1.55), Inches(11.9), [
        ["Dimension", "Conventional FDI / residual monitoring", "Conformal & runtime verification", "ASTRA's position"],
        ["Question asked", "Is a fault present now?", "Is this observation conformal to calibration?",
         "Where does the evidence stop being usable?"],
        ["Output", "binary flag per detector", "prediction set or alarm", "a per-stage evidence profile D_f"],
        ["Placement of monitors", "engineering convention or intuition", "wherever the score is defined",
         "proposed to follow measured observability"],
        ["Treatment of the estimator", "residuals treated as the signal", "estimator usually outside scope",
         "estimator treated as a transformation that can remove evidence"],
        ["Failure reporting", "missed detection / false alarm rates", "coverage guarantee",
         "adds: does the gate have the dynamic range to fire at all?"],
        ["What we add", "—", "—", "a measured separation of discriminability, calibration validity and operational stability"],
    ], widths=[2.0, 3.2, 3.2, 3.5], fs=10, hfs=9.5, rh=Inches(0.62))
    card(s, Inches(0.7), Inches(6.08), Inches(11.9), Inches(0.78),
         "The distinction the project has actually demonstrated",
         [("Statistical discriminability, calibration validity, and operational stability are three "
           "different properties. This system has the first without the third — measured, not argued.",
           INK, True)], TINT, DEEP, hs=11.5)

    # ---- 7 ROADMAP -------------------------------------------------------
    s = slide(prs, "Implementation roadmap and where we stand", "progress")
    stages = [
        ("E17", "Fault observability", "COMPLETE", 100),
        ("E17-P", "Corrected position injection", "COMPLETE", 100),
        ("E18", "OD-8 calibration", "COMPLETE", 100),
        ("E18-R1", "Run-local calibration", "COMPLETE", 100),
        ("E18-R2", "Matched-window calibration", "COMPLETE", 100),
        ("E18-R3", "Longer evaluation windows", "NEXT", 0),
        ("E19", "H7 monitor placement", "BLOCKED", 0),
        ("comma2k19", "Real-world data validation", "NOT STARTED", 0),
        ("highD", "Traffic-regime validation", "NOT STARTED", 0),
        ("E20", "Adversarial / lying sensor", "NOT STARTED", 0),
        ("CARLA", "Closed-loop validation", "NOT STARTED", 0),
    ]
    y = Inches(1.62)
    x0, xb = Inches(3.6), Inches(9.4)
    for code, name, st, pct in stages:
        col = CHIP[st]
        p(tb(s, Inches(0.7), y + Inches(0.03), Inches(1.15), Inches(0.28)), code, 11,
          INK, bold=True, first=True, space=0, font=HEAD)
        p(tb(s, Inches(1.85), y + Inches(0.04), Inches(1.7), Inches(0.28)), name, 10, BODY,
          first=True, space=0)
        rect(s, x0, y + Inches(0.02), xb - x0, Inches(0.24), fill=TINT2)
        if pct:
            rect(s, x0, y + Inches(0.02), Emu(int((xb - x0) * pct / 100)), Inches(0.24), fill=col)
        chip(s, Inches(9.62), y + Inches(0.02), st)
        y += Inches(0.34)
    rect(s, Inches(3.52), Inches(3.35), Emu(45720), Inches(0.42), fill=BAD)
    p(tb(s, Inches(3.66), Inches(3.42), Inches(3.0), Inches(0.3)),
      "◀  we are here", 11, BAD, bold=True, first=True, space=0)
    card(s, Inches(0.7), Inches(5.7), Inches(5.85), Inches(1.16),
         "Completed", [("5 of 11 stages. All internal validation. 4,290 closed-loop runs, "
                        "0 execution failures.", BODY)], TINT, OK, hs=11.5)
    card(s, Inches(6.75), Inches(5.7), Inches(5.85), Inches(1.16),
         "Blocking the rest",
         [("E19 needs an operational monitor. E18-R2 established that no policy currently "
           "has one. E18-R3 is the gate.", BODY)], TINT2, BAD, hs=11.5)
    foot(s, "Stages are gated: a stage is not entered because it exists in the plan. External datasets "
            "come after the internal causal chain, not before.")

    # ---- 8 CURRENT STATUS ------------------------------------------------
    s = slide(prs, "Experiment ledger", "current status")
    table(s, Inches(0.7), Inches(1.52), Inches(11.9), [
        ["Experiment", "Question", "Scale", "Verdict", "Outcome"],
        ["E17", "Stage-wise fault observability", "90 profiles · 2,160 runs",
         "COMPLETE", "Observability is heterogeneous; 1 of 6 faults well-posed"],
        ["E17-Position", "Does absorption survive correct injection?", "720 runs",
         "WITHDRAWN", "Absorbed in 0 of 12 cells — original result was artefact"],
        ["E18", "Can OD-8 be validly calibrated?", "72,000 ticks · 1,260 runs",
         "PARTIAL", "D_s does not predict detection; window defect found"],
        ["E18-R1", "Does run-local calibration recover P3?", "60 runs",
         "REJECTED", "P3 9/30; mechanism removed, estimator too noisy"],
        ["E18-R2", "Does matched-window calibration work?", "reused scores",
         "PARTIAL", "P1 13/30 — best of three schemes, still short"],
    ], widths=[1.5, 3.1, 2.0, 1.4, 4.2], fs=10, hfs=9.5, rh=Inches(0.58),
        colors={(1, 3): OK, (2, 3): BAD, (3, 3): WARN, (4, 3): BAD, (5, 3): WARN})
    for i, (v, l, c) in enumerate([("4,290", "closed-loop runs", DEEP),
                                   ("0", "execution failures", OK),
                                   ("6", "measurement defects found", WARN),
                                   ("5", "of those, self-caught", OK),
                                   ("0 / 30", "external validation [M-ext]", BAD)]):
        stat(s, Inches(0.7) + i * Inches(2.42), Inches(5.35), Inches(2.28), v, l, c, vs=30)
    foot(s, "Every invalidated record is retained and flagged, never deleted — an audit that erases its "
            "contaminated rows cannot be checked.")

    # ---- 9 RESULTS -------------------------------------------------------
    s = slide(prs, "What we have established", "results")
    left = [
        ("Fault observability is heterogeneous",
         "Four valid faults, four distinct behaviours: absorbed, absorbed-then-recovered, "
         "non-monotonic, persistent. A single absorption metric is ill-posed for five of six.",
         "ESTABLISHED"),
        ("One clean absorption result",
         "speed_stuck on policy P1: D 0.9625 → 0.5629 at the estimator, well-posed on 30/30 "
         "seeds, SD 0.0077.", "SUPPORTED"),
        ("A fault-injection channel was inert for 17 days",
         "POSITION_Y was regenerated from ground truth by the redundancy path. Demonstrated "
         "causally: estimator sees 0 with redundancy on, 3.18 with it off.", "ESTABLISHED"),
    ]
    right = [
        ("Discriminability does not predict detection",
         "17 of 28 cells disagree. D_L1 vs operational detection: Spearman ρ = −0.480, "
         "p = 0.0088 — the association runs backwards.", "ESTABLISHED"),
        ("Fault-induced alarm suppression",
         "11 of 28 cells, all p < 0.05. imu_dropout on P1 makes the monitor 55× LESS likely "
         "to alarm than clean operation — not better detection.", "SUPPORTED"),
        ("Two faults are undetectable at any severity",
         "speed_stuck and imu_dropout, correctly calibrated, on both valid policies. A property "
         "of the score, not of the threshold.", "ESTABLISHED"),
    ]
    for col, items in ((Inches(0.7), left), (Inches(6.75), right)):
        y = Inches(1.55)
        for h, body, st in items:
            rect(s, col, y, Inches(5.85), Inches(1.72), fill=TINT)
            chip(s, col + Inches(0.22), y + Inches(0.18), st)
            tf = tb(s, col + Inches(0.22), y + Inches(0.55), Inches(5.4), Inches(1.05))
            p(tf, h, 12, DEEP, bold=True, first=True, space=5, font=HEAD)
            p(tf, body, 10, BODY, space=0)
            y += Inches(1.84)
    foot(s, "All results are [M-syn] — synthetic plant, three policies, one severity per fault where "
            "parameterisable. Thirty seeds buys reproducibility, not external validity.")

    # ---- 10 THE CORE FINDING ---------------------------------------------
    s = slide(prs, "Why the monitor cannot be calibrated", "the central technical finding")
    p(tb(s, Inches(0.7), Inches(1.48), Inches(11.9), Inches(0.3)),
      "The threshold is in the right place. The score has no room around it.", 13, INK,
      bold=True, first=True, space=0, font=HEAD)
    for i, (v, l, c) in enumerate([("0.0097", "score standard deviation (P1)", DEEP),
                                   ("29 %", "of ticks within ±0.01 of the threshold", WARN),
                                   ("0.005", "threshold precision required", DEEP),
                                   ("0.013–0.028", "actual run-to-run baseline drift", BAD)]):
        stat(s, Inches(0.7) + i * Inches(3.03), Inches(1.95), Inches(2.85), v, l, c, vs=27)
    table(s, Inches(0.7), Inches(3.15), Inches(6.1), [
        ["Baseline shift", "P1 false-alarm rate"],
        ["−0.010", "0.20 %"],
        ["0 (calibrated)", "4.98 %"],
        ["+0.010", "29.17 %"],
        ["+0.020", "69.57 %"],
    ], widths=[2.4, 2.4], fs=11, colors={(2, 1): OK, (3, 1): BAD, (4, 1): BAD})
    card(s, Inches(7.05), Inches(3.15), Inches(5.55), Inches(1.94),
         "The problem in one sentence",
         [("The threshold must be placed to a precision of about 0.005 score units, while the "
           "score's own baseline wanders by 0.013–0.028 between runs.", INK, True),
          ("The precision required is 3–6× finer than the system's own reproducibility. "
           "Signal-to-noise at the decision boundary is below 1.", BODY)], TINT, BAD, hs=12)
    table(s, Inches(0.7), Inches(5.32), Inches(11.9), [
        ["Policy", "Overdispersion vs binomial", "Effective ticks (of 200)", "Alarm autocorrelation", "Mechanism"],
        ["P1", "4.2×", "11", "+0.359", "alarms cluster within a run"],
        ["P3", "9.2×", "2", "≈ 0", "baseline varies between runs"],
    ], widths=[1.1, 2.6, 2.4, 2.4, 3.4], fs=10.5)
    foot(s, "An ideal monitor with independent ticks would place 29.2 of 30 runs in band. Three calibration "
            "schemes were tried; the best result for any policy was 13/30. The limit is the score process.")

    # ---- 11 CLAIMS DISCIPLINE --------------------------------------------
    s = slide(prs, "What we refuse to claim", "claims discipline")
    p(tb(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.3)),
      "Six measurement defects have been found in this project's own code — five by self-audit, "
      "before publication. Each one withdrew a claim.", 11.5, BODY, first=True, space=0)
    withdrawn = [
        ("General L2a fault absorption", "position faults never reached the estimator"),
        ("Position-fault absorption", "0 of 12 cells after correct injection"),
        ("“Information is destroyed at L2a”", "a downstream statistic recovers part of it"),
        ("L6 detection-without-response gap", "the gate returns PASS on every tick"),
        ("H-regime operating-regime covariate", "Simpson's paradox: pooled −0.341, within-policy +0.836"),
        ("OD-8 provides operational monitoring", "no policy is calibrated on the evaluation window"),
    ]
    y = Inches(2.0)
    for h, why in withdrawn:
        rect(s, Inches(0.7), y, Inches(7.4), Inches(0.6), fill=TINT2)
        chip(s, Inches(0.9), y + Inches(0.17), "WITHDRAWN")
        p(tb(s, Inches(2.35), y + Inches(0.07), Inches(5.6), Inches(0.46)), h, 11, INK,
          bold=True, first=True, space=2)
        p(tb(s, Inches(2.35), y + Inches(0.30), Inches(5.6), Inches(0.26)), why, 9.5, BODY,
          first=True, space=0)
        y += Inches(0.68)
    card(s, Inches(8.35), Inches(2.0), Inches(4.25), Inches(1.95),
         "Terminology we enforce",
         [("D = 0.5 is chance-level separation, never “50 % accuracy”.", BODY),
          ("High D never implies a gate fired.", BODY),
          ("Alarm suppression is not better detection.", BODY)], TINT, DEEP, hs=11.5)
    card(s, Inches(8.35), Inches(4.15), Inches(4.25), Inches(1.93),
         "Why this is a strength",
         [("Every withdrawal came from auditing a result that looked too clean. The workflow now "
           "requires delivered-signal verification before any sweep is believed.", INK, True)],
         TINT2, OK, hs=11.5)
    foot(s, "Full status of every claim, with permitted wording, is maintained in CLAIM_LEDGER.md and "
            "updated after each experiment.")

    # ---- 12 NEXT STAGES --------------------------------------------------
    s = slide(prs, "Where this goes next", "further stages", dark=True)
    p(tb(s, Inches(0.7), Inches(1.42), Inches(11.9), Inches(0.32)),
      "The immediate question is narrow, cheap, and decisive.", 13,
      RGBColor(0x9F, 0xB8, 0xCC), first=True, space=0)
    steps = [
        ("E18-R3", "Longer evaluation windows",
         "At 11 effective ticks per 200, matching an independent monitor implies ~3,600 ticks. "
         "A pure compute change. Separates “not enough samples” from “wrong monitor”.", "NEXT"),
        ("E19 / H7", "Monitor placement",
         "Does the observability profile predict where a monitor works best? The prediction stays "
         "derived from D_s and is pre-registered — even though E18 suggests it will fail.", "BLOCKED"),
        ("E20", "Adversarial / lying sensor",
         "One channel lies while two stay honest, through the per-channel injection path already "
         "built. Not started; no detection claim is made.", "NOT STARTED"),
        ("External", "comma2k19 → highD → CARLA",
         "Real driving data, then traffic regimes, then closed loop. Gated behind a working "
         "internal measurement chain.", "NOT STARTED"),
    ]
    y = Inches(1.95)
    for code, name, body, st in steps:
        rect(s, Inches(0.7), y, Inches(11.9), Inches(1.03), fill=RGBColor(0x1A, 0x21, 0x4B))
        p(tb(s, Inches(1.0), y + Inches(0.17), Inches(1.6), Inches(0.3)), code, 13,
          RGBColor(0x6E, 0xC8, 0xE8), bold=True, first=True, space=0, font=HEAD)
        p(tb(s, Inches(1.0), y + Inches(0.52), Inches(2.1), Inches(0.32)), name, 10.5,
          RGBColor(0xA8, 0xC0, 0xD4), first=True, space=0)
        p(tb(s, Inches(3.3), y + Inches(0.22), Inches(7.6), Inches(0.7)), body, 10.5,
          RGBColor(0xC8, 0xD8, 0xE6), first=True, space=0)
        chip(s, Inches(11.15), y + Inches(0.38), st)
        y += Inches(1.13)
    card(s, Inches(0.7), Inches(6.52), Inches(11.9), Inches(0.62), "",
         [("If E18-R3 fails, the contribution reframes around the demonstrated limitation — a "
           "conformal score can separate faulted from clean runs at high AUC and still be unable to "
           "support a stable operational monitor. That is a real result.",
           RGBColor(0xD6, 0xE3, 0xEE), True)], RGBColor(0x2A, 0x33, 0x63), DEEP, hs=1)

    prs.save(str(out))
    print(f"  {len(prs.slides._sldIdLst)} slides -> {out}")


if __name__ == "__main__":
    build(Path(__file__).resolve().parent / "ASTRA_Project_Status_12.pptx")
